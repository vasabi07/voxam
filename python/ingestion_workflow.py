# NOTE: Unstructured is OPTIONAL - only used in legacy fallback methods
# Main pipeline uses lightweight libraries: PyMuPDF, python-docx, python-pptx
# Unstructured imports are done locally in _extract_pdf_legacy and _extract_auto

from langchain.chat_models import init_chat_model
from dotenv import load_dotenv
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from pydantic import BaseModel, Field
from typing import List, Literal, Optional, Tuple, Dict, Any
from enum import Enum
from openai import OpenAI
from neo4j import GraphDatabase
import os
import json
import time
import asyncio
from concurrent.futures import ThreadPoolExecutor
load_dotenv()

# OpenAI client for embeddings
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Embedding models
EMBED_MODEL_SMALL = "text-embedding-3-small"  # 1536 dims
EMBED_MODEL_LARGE = "text-embedding-3-large"  # 3072 dims (optional)

# Chunking parameters (optimized for educational content)
# Smaller chunks = better embedding precision, better RAG retrieval
CHUNK_MAX_CHARS = 4000           # Was 8000 - reduced for better precision
CHUNK_COMBINE_UNDER = 1500       # Was 3000 - don't over-combine small sections
CHUNK_NEW_AFTER = 3500           # Was 6000 - break earlier for better boundaries
CHUNK_OVERLAP = 500              # Was 200 - better context preservation at boundaries

# Neo4j credentials
NEO4J_URI = os.getenv("NEO4J_URI")
NEO4J_USER = os.getenv("NEO4J_USER")
NEO4J_PASSWORD = os.getenv("NEO4J_PASSWORD")

# Groq API for fast inference (5x faster than Cerebras)
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"

# Legacy API keys (fallback only)
DEEPINFRA_API_KEY = os.getenv("DEEPINFRA_API_KEY")
DEEPINFRA_URL = "https://api.deepinfra.com/v1/openai/chat/completions"
CEREBRAS_API_KEY = os.getenv("CEREBRAS_API_KEY")
CEREBRAS_URL = "https://api.cerebras.ai/v1/chat/completions"

# Google API for Gemini (handwriting OCR)
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")


# ============================================================
# Multi-Model OCR: Router & Extraction
# ============================================================

def detect_pdf_type(file_path: str, threshold: int = 200) -> str:
    """
    Detect if PDF is native (has text layer) or scanned/image-based.
    
    Returns:
        "native" - Has extractable text (use unstructured directly)
        "scanned" - No text layer, needs OCR
    """
    import fitz
    
    doc = fitz.open(file_path)
    total_chars = sum(len(page.get_text()) for page in doc)
    avg_chars_per_page = total_chars / len(doc) if len(doc) > 0 else 0
    doc.close()
    
    pdf_type = "native" if avg_chars_per_page > threshold else "scanned"
    print(f"📋 PDF type: {pdf_type} ({avg_chars_per_page:.0f} chars/page avg)")
    return pdf_type


def ocr_with_gemini(file_path: str) -> str:
    """
    Use Gemini 2.5 Flash for OCR - excellent for handwriting.
    Falls back to olmOCR if Gemini fails.
    """
    import httpx
    from pdf2image import convert_from_path
    from io import BytesIO
    import base64
    
    print("🖊️  Using Gemini 2.5 Flash for OCR (handwriting mode)...")
    
    if not GOOGLE_API_KEY:
        print("   ⚠️ No GOOGLE_API_KEY, falling back to olmOCR")
        return ocr_with_olmocr(file_path)
    
    try:
        # Convert PDF to images
        images = convert_from_path(file_path, dpi=150)
        
        all_text = []
        for page_num, img in enumerate(images, 1):
            # Convert to base64
            buffer = BytesIO()
            img.save(buffer, format="PNG")
            b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
            
            # Call Gemini API
            response = httpx.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GOOGLE_API_KEY}",
                headers={"Content-Type": "application/json"},
                json={
                    "contents": [{
                        "parts": [
                            {"text": "Extract ALL text from this image. Preserve formatting, line breaks, and structure. Just output the text, nothing else."},
                            {"inline_data": {"mime_type": "image/png", "data": b64}}
                        ]
                    }]
                },
                timeout=60.0
            )
            
            if response.status_code == 200:
                result = response.json()
                text = result["candidates"][0]["content"]["parts"][0]["text"]
                all_text.append(f"=== PAGE {page_num} ===\n{text}")
                print(f"   ✅ Page {page_num}: {len(text)} chars")
            else:
                print(f"   ⚠️ Page {page_num} failed: {response.status_code}")
                # Try olmOCR for this page
                page_text = _ocr_single_page_olmocr(img)
                all_text.append(f"=== PAGE {page_num} ===\n{page_text}")
        
        return "\n\n".join(all_text)
        
    except Exception as e:
        print(f"   ⚠️ Gemini failed: {e}, falling back to olmOCR")
        return ocr_with_olmocr(file_path)


def ocr_with_olmocr(file_path: str) -> str:
    """
    Use olmOCR via DeepInfra for OCR - good for printed scans.
    """
    import httpx
    from pdf2image import convert_from_path
    from io import BytesIO
    import base64
    
    print("📄 Using olmOCR for OCR (printed scan mode)...")
    
    if not DEEPINFRA_API_KEY:
        raise ValueError("DEEPINFRA_API_KEY not set")
    
    # Convert PDF to images
    images = convert_from_path(file_path, dpi=150)
    
    all_text = []
    for page_num, img in enumerate(images, 1):
        text = _ocr_single_page_olmocr(img)
        all_text.append(f"=== PAGE {page_num} ===\n{text}")
        print(f"   ✅ Page {page_num}: {len(text)} chars")
    
    return "\n\n".join(all_text)


def _ocr_single_page_olmocr(pil_image) -> str:
    """OCR a single PIL image using olmOCR via DeepInfra."""
    import httpx
    from io import BytesIO
    import base64
    
    # Convert to base64
    buffer = BytesIO()
    pil_image.save(buffer, format="PNG")
    b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
    
    response = httpx.post(
        DEEPINFRA_URL,
        headers={
            "Authorization": f"Bearer {DEEPINFRA_API_KEY}",
            "Content-Type": "application/json"
        },
        json={
            "model": "allenai/olmOCR-2-7B-1025",
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text from this image. Preserve structure and formatting."},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
                ]
            }],
            "max_tokens": 4000,
            "temperature": 0
        },
        timeout=90.0
    )
    
    if response.status_code == 200:
        result = response.json()
        return result["choices"][0]["message"]["content"]
    else:
        print(f"      ⚠️ olmOCR error: {response.status_code}")
        return ""


def embed_text(text: str, model: str = EMBED_MODEL_SMALL) -> List[float]:
    """
    Generate embeddings for text using OpenAI.
    Truncates long text to stay within token limits.
    """
    # Rough guardrail: ~8000 chars ≈ 2000 tokens (safe for text-embedding-3)
    # text = text[:8000]
    try:
        response = openai_client.embeddings.create(
            model=model,
            input=text
        )
        return response.data[0].embedding
    except Exception as e:
        print(f"❌ Failed to generate embedding: {e}")
        return []

def get_neo4j_driver():
    """Get Neo4j database driver"""
    if not all([NEO4J_URI, NEO4J_USER, NEO4J_PASSWORD]):
        raise ValueError("Missing Neo4j credentials in .env: NEO4J_URI, NEO4J_USER, NEO4J_PASSWORD")
    
    try:
        driver = GraphDatabase.driver(NEO4J_URI, auth=(NEO4J_USER, NEO4J_PASSWORD))
        # Test connection
        with driver.session() as session:
            session.run("RETURN 1")
        print("✅ Connected to Neo4j")
        return driver
    except Exception as e:
        print(f"❌ Cannot connect to Neo4j: {e}")
        raise e


"""
1.get the query
2. use unstructured.partition.pdf to extract content from the PDF
3. split text, images and tables
4. for each chunk create questions (must create a class to use structured output) and put all questions in a list
5. for each list item, create cyphers in neo4j graphdb 
6. for the actual chunk create a summary and have embeddings put that in the db as well
connections should be like user-> document-> topics -> chunks(embeddings) -> questions

"""

# ============================================================
# R2 Storage Configuration
# ============================================================

R2_ENDPOINT = os.getenv("R2_ENDPOINT")
R2_ACCESS_KEY_ID = os.getenv("R2_ACCESS_KEY_ID")
R2_SECRET_ACCESS_KEY = os.getenv("R2_SECRET_ACCESS_KEY")
R2_BUCKET = os.getenv("R2_BUCKET")
R2_PUBLIC_URL = os.getenv("R2_PUBLIC_URL", "https://r2.voxam.dev")


def get_r2_client():
    """Get boto3 S3 client configured for Cloudflare R2."""
    import boto3
    from botocore.client import Config

    if not all([R2_ENDPOINT, R2_ACCESS_KEY_ID, R2_SECRET_ACCESS_KEY, R2_BUCKET]):
        print("⚠️  R2 credentials not fully configured")
        return None

    return boto3.client(
        's3',
        endpoint_url=R2_ENDPOINT,
        aws_access_key_id=R2_ACCESS_KEY_ID,
        aws_secret_access_key=R2_SECRET_ACCESS_KEY,
        config=Config(signature_version='s3v4'),
        region_name='auto'
    )


def upload_to_r2(image_bytes: bytes, file_key: str, content_type: str = "image/png") -> Optional[str]:
    """
    Upload image to R2 and return public URL.

    Args:
        image_bytes: Raw image data
        file_key: R2 key (e.g., "documents/{doc_id}/images/img_001.png")
        content_type: MIME type of the image

    Returns:
        Public URL or None if upload failed
    """
    r2_client = get_r2_client()

    if not r2_client:
        print(f"   📤 [MOCK] Would upload {len(image_bytes)} bytes to R2: {file_key}")
        return f"{R2_PUBLIC_URL}/{file_key}"

    try:
        r2_client.put_object(
            Bucket=R2_BUCKET,
            Key=file_key,
            Body=image_bytes,
            ContentType=content_type
        )
        url = f"{R2_PUBLIC_URL}/{file_key}"
        print(f"   📤 Uploaded to R2: {file_key}")
        return url
    except Exception as e:
        print(f"   ❌ R2 upload failed: {e}")
        return None


# ============================================================
# PyMuPDF Text Extraction (Fast, No Unstructured)
# ============================================================

def extract_text_with_pymupdf(
    file_path: str,
    chunk_size: int = CHUNK_MAX_CHARS,
    chunk_overlap: int = CHUNK_OVERLAP
) -> List[Dict[str, Any]]:
    """
    Extract text from PDF using PyMuPDF (fast, no Unstructured dependency).
    Creates chunks based on pages with optional size limits.

    Args:
        file_path: Path to PDF file
        chunk_size: Target characters per chunk (default: CHUNK_MAX_CHARS)
        chunk_overlap: Overlap between chunks for context continuity (default: CHUNK_OVERLAP)

    Returns:
        List of chunk dicts with text, page_start, page_end
    """
    import fitz

    print(f"📄 Extracting text with PyMuPDF (fast mode)...")
    start_time = time.time()

    doc = fitz.open(file_path)
    page_count = len(doc)
    chunks = []
    current_chunk = ""
    current_page_start = 1
    total_chars = 0

    for page_num, page in enumerate(doc):
        page_text = page.get_text("text")
        total_chars += len(page_text)

        # Add page marker for reference
        page_text_with_marker = f"\n[PAGE {page_num + 1}]\n{page_text}"

        # Check if adding this page exceeds chunk size
        if len(current_chunk) + len(page_text_with_marker) > chunk_size and current_chunk:
            # Save current chunk
            chunks.append({
                "text": current_chunk.strip(),
                "page_start": current_page_start,
                "page_end": page_num,  # Previous page
                "chunk_index": len(chunks)
            })
            # Start new chunk with overlap
            overlap_text = current_chunk[-chunk_overlap:] if len(current_chunk) > chunk_overlap else ""
            current_chunk = overlap_text + page_text_with_marker
            current_page_start = page_num + 1
        else:
            current_chunk += page_text_with_marker

    # Don't forget the last chunk
    if current_chunk.strip():
        chunks.append({
            "text": current_chunk.strip(),
            "page_start": current_page_start,
            "page_end": page_count,
            "chunk_index": len(chunks)
        })

    doc.close()
    elapsed = time.time() - start_time
    print(f"   ✅ Extracted {total_chars:,} chars from {page_count} pages into {len(chunks)} chunks in {elapsed:.2f}s")
    return chunks


def extract_text_by_page(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from PDF, one chunk per page (simplest approach).
    Best for documents where pages are natural boundaries.

    Returns:
        List of chunk dicts, one per page
    """
    import fitz

    print(f"📄 Extracting text by page with PyMuPDF...")
    start_time = time.time()

    doc = fitz.open(file_path)
    chunks = []

    for page_num, page in enumerate(doc):
        page_text = page.get_text("text").strip()

        if page_text:  # Skip empty pages
            chunks.append({
                "text": page_text,
                "page_start": page_num + 1,
                "page_end": page_num + 1,
                "chunk_index": len(chunks)
            })

    doc.close()
    elapsed = time.time() - start_time
    print(f"   ✅ Extracted {len(chunks)} page chunks in {elapsed:.2f}s")
    return chunks


# ============================================================
# PyMuPDF Image Extraction
# ============================================================

def extract_images_from_pdf(
    file_path: str,
    min_width: int = 120,
    min_height: int = 120,
    min_bytes: int = 5000,
    max_aspect_ratio: float = 5.0,
    skip_header_footer_pct: float = 0.05
) -> Dict[int, List[dict]]:
    """
    Extract images from PDF using PyMuPDF with light filtering.
    Let vision LLM do the heavy lifting for educational vs decorative classification.

    Args:
        file_path: Path to PDF file
        min_width: Minimum image width in pixels (120 - only filter tiny icons)
        min_height: Minimum image height in pixels (120 - only filter tiny icons)
        min_bytes: Minimum file size in bytes (5KB - filter only very tiny images)
        max_aspect_ratio: Max width/height ratio (5:1 - filter extreme banners)
        skip_header_footer_pct: Skip images in top/bottom X% of page (0.05 = 5%)

    Returns:
        Dict mapping page_number (1-indexed) -> list of image dicts
    """
    import fitz
    from PIL import Image
    from io import BytesIO

    print(f"🖼️  Extracting images from PDF with PyMuPDF...")

    doc = fitz.open(file_path)
    images_by_page: Dict[int, List[dict]] = {}
    total_images = 0
    filtered_count = 0
    filter_reasons = {"size": 0, "dimensions": 0, "aspect": 0, "position": 0}

    for page_num, page in enumerate(doc):
        page_images = []
        image_list = page.get_images(full=True)
        page_height = page.rect.height

        for img_idx, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]

                # Filter 1: Minimum file size (skip tiny images)
                if len(image_bytes) < min_bytes:
                    filtered_count += 1
                    filter_reasons["size"] += 1
                    continue

                # Filter 2: Check actual dimensions
                width, height = None, None
                try:
                    pil_img = Image.open(BytesIO(image_bytes))
                    width, height = pil_img.size

                    # Skip small images (icons, bullets, small logos)
                    if width < min_width or height < min_height:
                        filtered_count += 1
                        filter_reasons["dimensions"] += 1
                        continue

                    # Skip very wide/narrow images (headers, banners, dividers)
                    aspect_ratio = max(width / height, height / width)
                    if aspect_ratio > max_aspect_ratio:
                        filtered_count += 1
                        filter_reasons["aspect"] += 1
                        continue

                except Exception:
                    # If we can't read dimensions, apply stricter size filter
                    if len(image_bytes) < 30000:
                        filtered_count += 1
                        filter_reasons["size"] += 1
                        continue

                # Get image position on page
                img_rects = page.get_image_rects(xref)
                bbox = img_rects[0] if img_rects else None
                y_pos = bbox.y0 if bbox else 0

                # Filter 3: Skip images in header/footer regions
                if bbox and page_height > 0:
                    y_center = (bbox.y0 + bbox.y1) / 2
                    relative_pos = y_center / page_height
                    # Skip if in top X% or bottom X% of page
                    if relative_pos < skip_header_footer_pct or relative_pos > (1 - skip_header_footer_pct):
                        filtered_count += 1
                        filter_reasons["position"] += 1
                        continue

                page_images.append({
                    "image_bytes": image_bytes,
                    "ext": image_ext,
                    "bbox": [bbox.x0, bbox.y0, bbox.x1, bbox.y1] if bbox else None,
                    "y_pos": y_pos,
                    "page": page_num + 1,
                    "width": width,
                    "height": height
                })
                total_images += 1

            except Exception as e:
                print(f"   ⚠️ Error extracting image {img_idx} from page {page_num + 1}: {e}")
                continue

        # Sort by y-position (reading order)
        if page_images:
            page_images.sort(key=lambda x: x["y_pos"])
            for idx, img in enumerate(page_images):
                img["index"] = idx
            images_by_page[page_num + 1] = page_images

    doc.close()
    print(f"   ✅ Extracted {total_images} images, filtered {filtered_count} (size:{filter_reasons['size']}, dim:{filter_reasons['dimensions']}, aspect:{filter_reasons['aspect']}, pos:{filter_reasons['position']})")
    return images_by_page


def classify_and_describe_image(
    image_bytes: bytes,
    page_num: int,
    vision_model: str = "gemini"  # Using Gemini Flash for higher rate limits
) -> Optional[Dict[str, Any]]:
    """
    Use vision LLM to classify image as educational or decorative, and describe if educational.
    Single API call does both filtering and description.

    Args:
        image_bytes: Raw image bytes
        page_num: Page number for context
        vision_model: Vision model to use (default: gemini)

    Returns:
        Dict with {"is_educational": bool, "description": str} or None on error
    """
    import base64
    import httpx
    import json

    # Encode image to base64
    b64_image = base64.b64encode(image_bytes).decode("utf-8")

    # Detect image type from bytes
    if image_bytes[:8] == b'\x89PNG\r\n\x1a\n':
        media_type = "image/png"
    elif image_bytes[:2] == b'\xff\xd8':
        media_type = "image/jpeg"
    else:
        media_type = "image/png"  # Default fallback

    prompt = """Analyze this image from an educational document (page """ + str(page_num) + """).

TASK: Determine if this is an EDUCATIONAL image or DECORATIVE/JUNK.

EDUCATIONAL images include:
- Diagrams explaining concepts or processes
- Charts, graphs showing data
- Figures illustrating scientific/technical concepts
- Flowcharts, mind maps
- Labeled illustrations (anatomy, circuits, etc.)
- Maps, timelines
- Mathematical figures, geometric shapes with labels

DECORATIVE/JUNK images include:
- Logos, brand marks, watermarks
- Page headers, footers, borders
- Decorative icons, bullet points
- Background patterns, textures
- Stock photos used for decoration
- Clipart without educational value
- Author photos, publisher logos

Respond in JSON format ONLY:
{
    "is_educational": true or false,
    "description": "Factual description of the image content (2-3 sentences). What does it show? What concepts does it illustrate?" OR null if decorative
}"""

    try:
        # Use Gemini Flash for higher rate limits
        response = httpx.post(
            f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GOOGLE_API_KEY}",
            headers={"Content-Type": "application/json"},
            json={
                "contents": [{
                    "parts": [
                        {"text": prompt},
                        {"inline_data": {"mime_type": media_type, "data": b64_image}}
                    ]
                }]
            },
            timeout=30.0
        )

        if response.status_code != 200:
            print(f"   ⚠️ Vision API error: {response.status_code}")
            return None

        result = response.json()
        content = result["candidates"][0]["content"]["parts"][0]["text"]

        # Parse JSON from response
        # Handle potential markdown code blocks
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0]
        elif "```" in content:
            content = content.split("```")[1].split("```")[0]

        parsed = json.loads(content.strip())
        return parsed

    except Exception as e:
        print(f"   ⚠️ Vision classification error: {e}")
        return None


def filter_and_describe_images(
    images_by_page: Dict[int, List[dict]],
    vision_model: str = "gemini",
    max_concurrent: int = 5,  # Gemini has higher rate limits
    max_retries: int = 2,
    fallback_on_failure: bool = True
) -> Dict[int, List[dict]]:
    """
    Apply vision LLM filtering to extracted images.
    Removes decorative images and adds descriptions to educational ones.

    Args:
        images_by_page: Output from extract_images_from_pdf()
        vision_model: Vision model for classification
        max_concurrent: Max parallel API calls
        max_retries: Retry count for failed API calls
        fallback_on_failure: If True, keep images that fail classification (assume educational)

    Returns:
        Filtered images_by_page with 'description' field added to educational images
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    total_images = sum(len(imgs) for imgs in images_by_page.values())
    if total_images == 0:
        return images_by_page

    print(f"🔍 Classifying {total_images} images with vision LLM...")

    # Flatten for parallel processing
    all_images = []
    for page_num, images in images_by_page.items():
        for img in images:
            all_images.append((page_num, img))

    educational_count = 0
    decorative_count = 0
    failed_count = 0

    def process_image_with_retry(page_num: int, img: dict) -> tuple:
        """Process image with retry logic."""
        last_error = None
        for attempt in range(max_retries + 1):
            try:
                result = classify_and_describe_image(
                    img["image_bytes"],
                    page_num,
                    vision_model
                )
                if result is not None:
                    return page_num, img, result, None
            except Exception as e:
                last_error = e
                if attempt < max_retries:
                    time.sleep(1 * (attempt + 1))  # Exponential backoff
        return page_num, img, None, last_error

    # Process in parallel
    with ThreadPoolExecutor(max_workers=max_concurrent) as executor:
        futures = [
            executor.submit(process_image_with_retry, page_num, img)
            for page_num, img in all_images
        ]

        for future in as_completed(futures):
            try:
                page_num, img, result, error = future.result()

                if error:
                    failed_count += 1
                    if fallback_on_failure:
                        # Assume educational if we can't classify - better to keep than lose
                        img["description"] = f"Image from page {page_num} (classification unavailable)"
                        img["is_educational"] = True
                        img["classification_failed"] = True
                        educational_count += 1
                    else:
                        img["is_educational"] = False
                elif result and result.get("is_educational"):
                    img["description"] = result.get("description", "Educational image")
                    img["is_educational"] = True
                    educational_count += 1
                else:
                    img["is_educational"] = False
                    decorative_count += 1
            except Exception as e:
                # Future itself failed
                failed_count += 1
                print(f"   ⚠️ Unexpected error processing image: {e}")

    # Filter out decorative images
    filtered_by_page = {}
    for page_num, images in images_by_page.items():
        educational = [img for img in images if img.get("is_educational", False)]
        if educational:
            # Re-index after filtering
            for idx, img in enumerate(educational):
                img["index"] = idx
            filtered_by_page[page_num] = educational

    status_parts = [f"{educational_count} educational", f"{decorative_count} decorative removed"]
    if failed_count > 0:
        status_parts.append(f"{failed_count} failed (kept as fallback)" if fallback_on_failure else f"{failed_count} failed")
    print(f"   ✅ Vision filter: {', '.join(status_parts)}")

    return filtered_by_page


def upload_images_to_r2(
    images_by_page: Dict[int, List[dict]],
    doc_id: str,
    max_concurrent: int = 10
) -> Dict[str, str]:
    """
    Upload all extracted images to R2 in parallel.

    Args:
        images_by_page: Dict mapping page_num -> list of image dicts
        doc_id: Document ID for R2 path
        max_concurrent: Max parallel uploads (default 10)

    Returns:
        Dict mapping "page_N_img_M" -> image_url
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    print(f"📤 Uploading images to R2...")

    # Collect all upload tasks
    upload_tasks = []
    for page_num in sorted(images_by_page.keys()):
        for img in images_by_page[page_num]:
            img_idx = img["index"]
            ext = img.get("ext", "png")
            file_key = f"documents/{doc_id}/images/p{page_num}_img{img_idx}.{ext}"
            content_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"
            upload_tasks.append((page_num, img_idx, img["image_bytes"], file_key, content_type))

    if not upload_tasks:
        return {}

    # Upload in parallel
    def do_upload(task):
        page_num, img_idx, image_bytes, file_key, content_type = task
        url = upload_to_r2(image_bytes, file_key, content_type)
        if url:
            print(f"   📤 Uploaded to R2: {file_key}")
            return f"page_{page_num}_img_{img_idx}", url
        return None, None

    image_index: Dict[str, str] = {}
    with ThreadPoolExecutor(max_workers=max_concurrent) as executor:
        futures = [executor.submit(do_upload, task) for task in upload_tasks]
        for future in as_completed(futures):
            key, url = future.result()
            if key and url:
                image_index[key] = url

    print(f"   ✅ Uploaded {len(image_index)} images to R2")
    return image_index


def describe_image(image_bytes: bytes, vision_llm=None) -> str:
    """
    Generate a factual description of an image using Gemini Flash.
    Focuses on WHAT is shown, not interpretation.

    Note: vision_llm parameter is kept for backward compatibility but ignored.
    We use Gemini Flash directly for higher rate limits.
    """
    import base64

    b64 = base64.b64encode(image_bytes).decode("utf-8")

    prompt = """Describe what this image shows factually. Include:
- Type of visual (diagram, graph, flowchart, table, photo, illustration)
- All visible labels, text, and numbers
- Spatial relationships (what connects to what)
- Data values if present

Be specific and factual. Do NOT explain concepts or summarize - just describe what is visible.

Example: "Diagram of a neuron showing: cell body labeled 'soma' in center, dendrites branching left, single axon extending right with myelin sheath segments (labeled), terminal buttons at end. Arrow indicates signal direction from dendrites to axon terminal."
"""

    try:
        response = httpx.post(
            f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GOOGLE_API_KEY}",
            headers={"Content-Type": "application/json"},
            json={
                "contents": [{
                    "parts": [
                        {"text": prompt},
                        {"inline_data": {"mime_type": "image/png", "data": b64}}
                    ]
                }]
            },
            timeout=30.0
        )

        if response.status_code == 200:
            result = response.json()
            return result["candidates"][0]["content"]["parts"][0]["text"]
        else:
            print(f"   ⚠️ Gemini error: {response.status_code}")
            return "Image description unavailable"

    except Exception as e:
        print(f"   ⚠️ Image description failed: {e}")
        return "Image description unavailable"


# ============================================================
# LLM Hierarchy Creation
# ============================================================

def create_hierarchy_with_llm(
    chunks: List['ContentBlock'],
    model: str = "llama-3.3-70b"
) -> dict:
    """
    Use LLM to create chapter/section hierarchy from chunk content.
    Analyzes first lines and key content to determine document structure.

    Returns:
        {
            "chapters": [
                {
                    "title": "Chapter 1: ...",
                    "sections": [
                        {"title": "1.1 ...", "chunk_indices": [0, 1, 2]}
                    ]
                }
            ]
        }
    """
    import httpx

    print("🏗️  Creating hierarchy with LLM...")

    # Build informative input - first lines + key content hints
    chunk_summaries = []
    for i, chunk in enumerate(chunks):
        text = chunk.text_content if hasattr(chunk, 'text_content') else chunk.get('text', '')
        page = chunk.page_number if hasattr(chunk, 'page_number') else chunk.get('page_start', i+1)

        # Get first 3 meaningful lines for context
        lines = [l.strip() for l in text.split('\n') if l.strip() and len(l.strip()) > 10]
        preview = ' | '.join(lines[:3])[:200] if lines else "Empty chunk"

        chunk_summaries.append(f"Chunk {i} [Page {page}]: {preview}")

    prompt = f"""You are analyzing an educational document to create a chapter/section hierarchy.

DOCUMENT CHUNKS:
{chr(10).join(chunk_summaries)}

YOUR TASK: Create a logical hierarchy that organizes these chunks into chapters and sections.

ANALYSIS GUIDELINES:
1. **Detect existing structure**: Look for chapter/unit markers like "Chapter 1", "Unit 1", "Module 1", "Part I", etc.
2. **Identify topic boundaries**: When chunks shift to a new major topic, that's likely a new chapter
3. **Group related content**: Chunks discussing the same concept should be in the same section
4. **Preserve reading order**: Chunk indices must be in ascending order within sections
5. **Create meaningful titles**: Use actual headings from text, or create descriptive titles based on content

RULES:
- Every chunk index (0 to {len(chunks)-1}) must appear exactly ONCE
- Chunk indices within a section must be consecutive (e.g., [0,1,2] not [0,2,4])
- Aim for 2-6 chapters with 1-4 sections each
- Section titles should be specific and descriptive

OUTPUT FORMAT - Return ONLY valid JSON:
{{
  "chapters": [
    {{
      "title": "Chapter 1: [Descriptive Title]",
      "sections": [
        {{"title": "[Section Title]", "chunk_indices": [0, 1]}},
        {{"title": "[Section Title]", "chunk_indices": [2]}}
      ]
    }},
    {{
      "title": "Chapter 2: [Descriptive Title]",
      "sections": [
        {{"title": "[Section Title]", "chunk_indices": [3, 4, 5]}}
      ]
    }}
  ]
}}

IMPORTANT: Output ONLY the JSON, no explanation or markdown formatting."""

    try:
        # Use Groq Llama 3.1 8B for hierarchy (fast, cheap, valid JSON output)
        api_url = GROQ_URL if GROQ_API_KEY else (CEREBRAS_URL if CEREBRAS_API_KEY else DEEPINFRA_URL)
        api_key = GROQ_API_KEY if GROQ_API_KEY else (CEREBRAS_API_KEY if CEREBRAS_API_KEY else DEEPINFRA_API_KEY)
        model = "llama-3.1-8b-instant" if GROQ_API_KEY else ("llama-3.3-70b" if CEREBRAS_API_KEY else "meta-llama/Llama-3.3-70B-Instruct")

        response = httpx.post(
            api_url,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": model,
                "messages": [
                    {"role": "system", "content": "You are a document structure analyzer. Output valid JSON only."},
                    {"role": "user", "content": prompt}
                ],
                "max_tokens": 2000,
                "temperature": 0
            },
            timeout=30.0  # Groq is very fast
        )

        if response.status_code != 200:
            print(f"   ⚠️ LLM hierarchy error: {response.status_code}")
            return _fallback_hierarchy(chunks)

        result = response.json()
        text = result["choices"][0]["message"]["content"]

        # Parse JSON
        start = text.find('{')
        end = text.rfind('}') + 1
        if start != -1 and end > start:
            hierarchy = json.loads(text[start:end])

            # Validate structure
            if "chapters" in hierarchy and len(hierarchy["chapters"]) > 0:
                total_chunks = sum(
                    len(sec.get("chunk_indices", []))
                    for ch in hierarchy["chapters"]
                    for sec in ch.get("sections", [])
                )
                print(f"   ✅ Created hierarchy: {len(hierarchy['chapters'])} chapters, {total_chunks} chunks mapped")
                return hierarchy

        print("   ⚠️ Invalid hierarchy JSON, using fallback")
        return _fallback_hierarchy(chunks)

    except Exception as e:
        print(f"   ⚠️ Hierarchy creation failed: {e}")
        return _fallback_hierarchy(chunks)


def _fallback_hierarchy(chunks: List['ContentBlock']) -> dict:
    """
    Fallback: Create simple hierarchy based on page groupings.
    Groups every 3-5 chunks into a section.
    """
    print("   📋 Using fallback page-based hierarchy")

    sections = []
    chunk_count = len(chunks)
    chunks_per_section = max(2, min(5, chunk_count // 5))

    for i in range(0, chunk_count, chunks_per_section):
        end_idx = min(i + chunks_per_section, chunk_count)
        indices = list(range(i, end_idx))

        # Use first chunk's first line as section title
        first_line = chunks[i].text_content.split('\n')[0][:60] if chunks[i].text_content else f"Section {i//chunks_per_section + 1}"

        sections.append({
            "title": first_line,
            "chunk_indices": indices
        })

    return {
        "chapters": [{
            "title": "Document Content",
            "sections": sections
        }]
    }


def apply_hierarchy_to_chunks(
    chunks: List['ContentBlock'],
    hierarchy: dict
) -> List['ContentBlock']:
    """
    Apply hierarchy metadata to chunks.
    """
    for chapter in hierarchy.get("chapters", []):
        chapter_title = chapter.get("title", "Unknown Chapter")

        for section in chapter.get("sections", []):
            section_title = section.get("title", "Unknown Section")

            for idx in section.get("chunk_indices", []):
                if 0 <= idx < len(chunks):
                    chunks[idx].chapter_title = chapter_title
                    chunks[idx].section_title = section_title
                    chunks[idx].heading_level = 2  # Section level

    return chunks


def create_topic_content_blocks(
    raw_chunks: List[Dict],
    hierarchy: dict
) -> List['ContentBlock']:
    """
    Combine raw chunks into topic-level ContentBlocks based on LLM hierarchy.

    Each section in the hierarchy becomes a single ContentBlock with:
    - Combined text from all chunks in that section
    - Page range spanning all included chunks
    - Chapter/section metadata

    Args:
        raw_chunks: List of chunk dicts from chunk_pages() with keys:
                   text, index, page_start, page_end
        hierarchy: LLM-generated hierarchy with chapters/sections/chunk_indices

    Returns:
        List of ContentBlocks, one per section (topic)
    """
    content_blocks = []
    block_index = 0

    for chapter in hierarchy.get("chapters", []):
        chapter_title = chapter.get("title", "Unknown Chapter")

        for section in chapter.get("sections", []):
            section_title = section.get("title", "Unknown Section")
            chunk_indices = section.get("chunk_indices", [])

            if not chunk_indices:
                continue

            # Validate indices
            valid_indices = [i for i in chunk_indices if 0 <= i < len(raw_chunks)]
            if not valid_indices:
                continue

            # Combine text from all chunks in this section
            section_texts = []
            page_starts = []
            page_ends = []

            for idx in valid_indices:
                chunk = raw_chunks[idx]
                section_texts.append(chunk["text"])
                page_starts.append(chunk.get("page_start", 1))
                page_ends.append(chunk.get("page_end", 1))

            # Create the topic-level ContentBlock
            block = ContentBlock()
            block.chunk_index = block_index
            block.text_content = "\n\n".join(section_texts)
            block.combined_context = block.text_content  # Will be enriched later

            # Page range spans all included chunks
            block.page_start = min(page_starts)
            block.page_end = max(page_ends)
            block.page_number = block.page_start  # Primary page reference

            # Hierarchy metadata
            block.chapter_title = chapter_title
            block.section_title = section_title
            block.heading_level = 2  # Section level

            content_blocks.append(block)
            block_index += 1

    print(f"   ✅ Created {len(content_blocks)} topic-level ContentBlocks from {len(raw_chunks)} raw chunks")
    return content_blocks


def match_images_to_chunks(
    chunks: List['ContentBlock'],
    images_by_page: Dict[int, List[dict]],
    image_index: Dict[str, str],
    vision_llm=None
) -> List['ContentBlock']:
    """
    Match extracted images to content chunks based on page numbers.
    Uses descriptions from vision filtering if available.

    Args:
        chunks: List of ContentBlocks with page_number set
        images_by_page: Dict from extract_images_from_pdf() - may have 'description' from vision filter
        image_index: Dict from upload_images_to_r2() mapping key -> URL
        vision_llm: Optional LLM for image descriptions (legacy, prefer vision filtering)

    Returns:
        Updated chunks with image_urls, image_descriptions, and figure_map populated
    """
    print("🔗 Matching images to content chunks...")

    images_matched = 0

    for chunk in chunks:
        # Get page range for this chunk
        page_start = getattr(chunk, 'page_start', chunk.page_number)
        page_end = getattr(chunk, 'page_end', chunk.page_number)

        # Initialize figure tracking for this chunk
        chunk.figure_map = {}  # Maps figure_num -> {"url": ..., "description": ...}
        figure_num = 1

        # Collect images from all pages in chunk's range
        for page_num in range(page_start, page_end + 1):
            page_images = images_by_page.get(page_num, [])

            for img in page_images:
                img_idx = img.get("index", 0)
                key = f"page_{page_num}_img_{img_idx}"

                if key in image_index:
                    url = image_index[key]
                    chunk.image_urls.append(url)

                    # Use description from vision filtering if available
                    if img.get("description"):
                        desc = img["description"]
                    elif vision_llm:
                        desc = describe_image(img["image_bytes"], vision_llm)
                    else:
                        desc = f"Image from page {page_num}"

                    chunk.image_descriptions.append(desc)

                    # Build figure map for QP generation
                    chunk.figure_map[figure_num] = {
                        "url": url,
                        "description": desc,
                        "page": page_num
                    }
                    figure_num += 1
                    images_matched += 1

    print(f"   ✅ Matched {images_matched} images to {len(chunks)} chunks")
    return chunks


def build_combined_context_with_figures(block: 'ContentBlock') -> str:
    """
    Build combined_context including text and numbered figure descriptions.
    This context is used for embeddings and QP generation.

    Args:
        block: ContentBlock with text_content and figure_map

    Returns:
        Combined context string with [FIGURES] section
    """
    parts = [block.text_content]

    # Add figure descriptions if any
    if hasattr(block, 'figure_map') and block.figure_map:
        figure_lines = ["\n[FIGURES IN THIS SECTION]"]
        for fig_num, fig_data in sorted(block.figure_map.items()):
            desc = fig_data.get("description", "")
            page = fig_data.get("page", "?")
            figure_lines.append(f"- Figure {fig_num} (p.{page}): {desc}")
        parts.append("\n".join(figure_lines))

    # Add any other captions (from Unstructured, if used)
    if block.image_captions:
        parts.append("\n[IMAGE CAPTIONS]\n" + "\n".join(block.image_captions))

    if block.table_descriptions:
        parts.append("\n[TABLES]\n" + "\n".join(block.table_descriptions))

    return "\n\n".join(parts)


def add_image_context_to_questions(
    block: 'ContentBlock'
) -> List['Question']:
    """
    Add image URLs and descriptions to questions generated for a block.
    Questions that reference images will have the image_url field populated.
    """
    if not block.questions or not block.image_urls:
        return block.questions

    # For simplicity, attach first image to all questions from this block
    # More sophisticated matching could parse question text for image references
    primary_image_url = block.image_urls[0] if block.image_urls else None
    primary_image_desc = block.image_descriptions[0] if block.image_descriptions else None

    for question in block.questions:
        # Check if question likely references an image
        q_text_lower = question.text.lower()
        image_keywords = ["diagram", "figure", "image", "graph", "chart", "shown", "illustrated", "picture"]

        if any(keyword in q_text_lower for keyword in image_keywords):
            question.image_url = primary_image_url
            question.image_description = primary_image_desc

    return block.questions


class BloomLevel(str, Enum):
    REMEMBER = "remember"
    UNDERSTAND = "understand" 
    APPLY = "apply"
    ANALYZE = "analyze"
    EVALUATE = "evaluate"
    CREATE = "create"

class Difficulty(str, Enum):
    BASIC = "basic"
    INTERMEDIATE = "intermediate"
    ADVANCED = "advanced"

class QuestionType(str, Enum):
    LONG_ANSWER = "long_answer"
    MULTIPLE_CHOICE = "multiple_choice"

class Question(BaseModel):
    """
    Question model - simplified for efficiency.

    key_points is the core grading data. The correction_agent generates
    explanations, model answers, and rubrics at runtime from key_points + content_block.

    spoken_text and spoken_options are pre-computed because questions are read verbatim
    (not rephrased by LLM like content blocks).
    """
    # Core question data
    text: str = Field(description="The question text")
    bloom_level: BloomLevel = Field(description="Bloom's taxonomy level")
    difficulty: Difficulty = Field(description="Question difficulty level")
    question_type: QuestionType = Field(description="Question type (mcq or long_answer)")
    expected_time: int = Field(description="Expected time to answer in seconds")

    # THE CORE GRADING DATA - correction_agent uses this
    key_points: List[str] = Field(description="Key points the answer should cover (2-4 items)")

    # MCQ only
    options: Optional[List[str]] = Field(default=None, description="MCQ options [A, B, C, D]")
    correct_answer: Optional[str] = Field(default=None, description="Correct option letter (A/B/C/D)")

    # TTS support (questions are read verbatim, not rephrased by LLM)
    spoken_text: Optional[str] = Field(default=None, description="TTS-optimized question text (equations as words)")
    spoken_options: Optional[List[str]] = Field(default=None, description="TTS-optimized options for MCQ")

    # Image/figure fields for multimodal questions
    figure_ref: Optional[int] = Field(default=None, description="Figure number referenced (1, 2, etc.) - links to figure_map")
    image_url: Optional[str] = Field(default=None, description="R2 URL for question image (populated from figure_ref)")
    image_description: Optional[str] = Field(default=None, description="Factual description of image for LLM context")

class QuestionSet(BaseModel):
    """Container for all questions generated from a content block"""
    long_answer_questions: List[Question] = Field(description="Long answer questions")
    multiple_choice_questions: List[Question] = Field(description="Multiple choice questions") 
    total_questions: int = Field(description="Total number of questions")

class ContentBlock:
    """
    Enhanced content block with rich metadata for educational content.

    NOTE: No spoken_content field here - TTS preprocessing happens at RUNTIME
    on LLM output in realtime_base.py (LLM rephrases content anyway).
    """
    def __init__(self):
        # Core content
        self.chunk_index: int = 0
        self.text_content: str = ""
        self.combined_context: str = ""

        # Legacy fields (still used)
        self.related_tables: List[str] = []
        self.image_captions: List[str] = []
        self.table_descriptions: List[str] = []

        # Questions and embeddings
        self.questions: List[Question] = []
        self.embeddings: List[float] = []
        self.meta: dict = {}

        # Page tracking
        self.page_number: int = 1
        self.page_start: int = 1
        self.page_end: int = 1
        self.bbox: Optional[List[float]] = None  # [x1, y1, x2, y2] for scroll-to

        # Hierarchy (populated by LLM structuring phase)
        self.chapter_title: Optional[str] = None
        self.section_title: Optional[str] = None
        self.heading_level: int = 0  # 1=chapter, 2=section, 3=subsection

        # Image fields
        self.image_urls: List[str] = []
        self.image_descriptions: List[str] = []
        self.figure_map: Dict[int, dict] = {}  # {figure_num: {"url", "description", "page"}}

        # NEW: Content type classification
        # Values: definition, example, theorem, proof, procedure, code, equation, narrative
        self.content_type: str = "narrative"

        # NEW: Extracted structured content
        self.definitions: List[Dict[str, str]] = []  # [{"term": ..., "definition": ...}]
        self.procedure_steps: List[str] = []  # ["Step 1: ...", "Step 2: ..."]
        self.equations: List[Dict[str, str]] = []  # [{"latex": ..., "spoken": ..., "context": ...}]
        self.code_blocks: List[Dict[str, str]] = []  # [{"language": ..., "code": ..., "description": ...}]

        # NEW: Enhanced tables (structured, not just descriptions)
        self.tables: List[Dict] = []  # [{"headers": [], "rows": [[]], "description": ..., "spoken": ...}]

        # NEW: Image type classification
        self.image_types: List[str] = []  # ["circuit", "anatomy", "graph", "flowchart", ...]

class IngestionPipeline:
    def __init__(self, config):
        self.config = config
        self.vision_llm = init_chat_model(model=config.get("vision_llm", "gpt-4o-mini"), temperature=0)
        self.text_llm = init_chat_model(model=config.get("text_llm", "gpt-4.1"), temperature=0)
        self.neo4j_driver = None
        
        # Initialize Neo4j if credentials available
        try:
            self.neo4j_driver = get_neo4j_driver()
        except Exception as e:
            print(f"⚠️  Neo4j not available: {e}")
            print("Continuing without graph database persistence...")


    def extract_document(self, file_path: str) -> List[ContentBlock]:
        """Extract content from any supported document format (PDF, DOCX, PPTX, MD, images)"""
        import os
        ext = os.path.splitext(file_path)[1].lower()

        # Image formats (OCR required)
        image_formats = ['.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.webp']
        # Document formats
        doc_formats = ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.md', '.markdown']
        supported_formats = doc_formats + image_formats

        if ext not in supported_formats:
            raise ValueError(f"Unsupported file format: {ext}. Supported: {supported_formats}")

        print(f"📄 Detected format: {ext}")

        # Image files → OCR extraction
        if ext in image_formats:
            return self._extract_image(file_path)
        # Use format-specific partitioner for best results, fallback to auto
        elif ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx(file_path)
        elif ext in ['.md', '.markdown']:
            return self._extract_markdown(file_path)
        else:
            return self._extract_auto(file_path)
    
    def _process_text_to_blocks(
        self,
        pages_text: List[str],
        format_name: str = "Document",
        start_time: float = None
    ) -> List[ContentBlock]:
        """
        UNIFIED PIPELINE: Convert extracted text pages into topic-level ContentBlocks.

        This is the common path for ALL document formats after initial text extraction:
        1. Chunking with smart boundaries
        2. LLM hierarchy creation (chapters/sections)
        3. Topic-level content block grouping
        4. Content type detection and enrichment

        Args:
            pages_text: List of text strings (one per page/slide)
            format_name: Format name for logging (PDF, DOCX, PPTX, Image, etc.)
            start_time: Start time for elapsed calculation

        Returns:
            List of enriched, topic-level ContentBlocks
        """
        from lib.text_chunker import chunk_pages
        from lib.content_detector import (
            detect_content_type,
            extract_definitions,
            extract_procedure_steps,
        )
        from lib.math_to_speech import extract_equations

        if start_time is None:
            start_time = time.time()

        # Step 1: Chunking
        print(f"📄 Step 1: Chunking {len(pages_text)} pages...")
        raw_chunks = chunk_pages(pages_text, max_chars=4000, min_chars=500, overlap=200)
        print(f"   ✅ Created {len(raw_chunks)} chunks")

        if not raw_chunks:
            print("   ⚠️ No chunks created")
            return []

        # Step 2: LLM Hierarchy Creation
        print("🏗️  Step 2: Creating hierarchy with LLM...")
        temp_blocks = []
        for chunk in raw_chunks:
            temp = ContentBlock()
            temp.text_content = chunk["text"]
            temp.page_number = chunk["page_start"]
            temp.page_start = chunk["page_start"]
            temp.page_end = chunk["page_end"]
            temp_blocks.append(temp)

        hierarchy = create_hierarchy_with_llm(temp_blocks)

        # Step 3: Topic-level Content Blocks
        print("📦 Step 3: Creating topic-level ContentBlocks...")
        content_blocks = create_topic_content_blocks(raw_chunks, hierarchy)

        # Step 4: Content Detection and Enrichment
        print("🔍 Step 4: Content analysis and enrichment...")
        for block in content_blocks:
            # Content type classification
            block.content_type = detect_content_type(block.text_content).value

            # Extract definitions
            if block.content_type == "definition" or "definition" in block.text_content.lower():
                block.definitions = extract_definitions(block.text_content)

            # Extract procedure steps
            if block.content_type == "procedure":
                block.procedure_steps = extract_procedure_steps(block.text_content)

            # Extract equations with spoken forms
            block.equations = extract_equations(block.text_content)

            # Extract code blocks
            block.code_blocks = self._extract_code_blocks(block.text_content)

        # Build combined context
        for block in content_blocks:
            block.combined_context = self._combine_context(block)

        elapsed = time.time() - start_time
        print(f"✅ {format_name} processing completed in {elapsed:.2f}s ({len(content_blocks)} topic blocks)")
        return content_blocks

    def _extract_pdf(self, pdf_path: str) -> List[ContentBlock]:
        """
        Extract PDF content using fast PyMuPDF + OCR fallback, organized into topic-level ContentBlocks.

        Pipeline:
        1. PyMuPDF text extraction
        2. Encoding check + OCR fallback for problem pages
        3. → Unified pipeline (chunking, hierarchy, topic blocks, enrichment)
        """
        from lib.encoding_check import extract_pdf_with_fallback

        print("⏳ Starting PDF extraction (fast PyMuPDF mode)...")
        start_time = time.time()

        # Step 1: Fast extraction with OCR fallback for problem pages
        print("📄 PyMuPDF extraction with encoding check + OCR fallback...")
        pages_text = extract_pdf_with_fallback(pdf_path, doc_id=None)
        extract_time = time.time() - start_time
        print(f"   ✅ Extracted {len(pages_text)} pages in {extract_time:.2f}s")

        if not pages_text or all(not p.strip() for p in pages_text):
            print("   ⚠️ No text extracted, falling back to full OCR")
            return self._extract_scanned_pdf(pdf_path, start_time)

        # Step 2: Unified pipeline
        return self._process_text_to_blocks(pages_text, "PDF", start_time)

    def _extract_pdf_legacy(self, pdf_path: str) -> List[ContentBlock]:
        """
        Legacy PDF extraction using Unstructured hi_res strategy.
        Kept as fallback if fast extraction produces poor results.

        NOTE: Requires Unstructured + Tesseract to be installed.
        """
        # Import Unstructured locally (optional dependency)
        try:
            from unstructured.partition.pdf import partition_pdf
        except ImportError:
            raise ImportError(
                "Unstructured is required for legacy PDF extraction. "
                "Install with: pip install unstructured[pdf]"
            )

        print("⏳ Starting PDF extraction (legacy hi_res mode)...")
        start_time = time.time()

        # Router: Check if native or scanned PDF
        pdf_type = detect_pdf_type(pdf_path)

        if pdf_type == "scanned":
            # Scanned PDF: OCR first, then partition text
            return self._extract_scanned_pdf(pdf_path, start_time)

        # Native PDF: Use unstructured directly
        chunks = partition_pdf(
            filename=pdf_path,
            infer_table_structure=True,
            strategy="hi_res",
            extract_image_block_types=["Image"],
            extract_image_block_to_payload=True,
            chunking_strategy="by_title",
            max_characters=CHUNK_MAX_CHARS,
            combine_text_under_n_chars=CHUNK_COMBINE_UNDER,
            new_after_n_chars=CHUNK_NEW_AFTER,
            overlap=CHUNK_OVERLAP,
        )


        parse_time = time.time() - start_time
        print(f"   ✅ PDF parsed in {parse_time:.2f}s")

        # Collect all images with their block indices for batch processing
        all_images: List[Tuple[int, str]] = []  # (block_idx, base64)

        content_blocks = []
        for idx, chunk in enumerate(chunks):
            block = ContentBlock()
            block.chunk_index = idx

            # All chunks get text content
            block.text_content = str(chunk)

            # Extract page number (primary) and optional bbox (nice-to-have)
            if hasattr(chunk, "metadata"):
                block.page_number = chunk.metadata.page_number if chunk.metadata.page_number else 1

                # Optional bbox for scroll-to position
                if hasattr(chunk.metadata, "coordinates") and chunk.metadata.coordinates:
                    coords = chunk.metadata.coordinates
                    if coords.points:
                        points = list(coords.points)
                        xs = [p[0] for p in points]
                        ys = [p[1] for p in points]
                        block.bbox = [min(xs), min(ys), max(xs), max(ys)]

            # Handle different chunk types - collect images but don't caption yet
            if "CompositeElement" in str(type(chunk)):
                # Extract image base64 data (don't caption yet)
                images = self._extract_image_base64_from_chunk(chunk)
                for img_b64 in images:
                    all_images.append((idx, img_b64))
                # CompositeElements might also contain tables
                block.related_tables.extend(self._extract_tables_from_chunk(chunk))

            elif "Table" in str(type(chunk)):
                block.related_tables = [chunk]

            content_blocks.append(block)

        # Batch caption all images in parallel (returns tuples of caption, image_type)
        if all_images:
            print(f"🖼️  Captioning {len(all_images)} images in parallel...")
            caption_start = time.time()
            caption_results = self._batch_caption_images([img for _, img in all_images])

            # Assign captions and image types back to their blocks
            for (block_idx, _), (caption, image_type) in zip(all_images, caption_results):
                content_blocks[block_idx].image_captions.append(caption)
                content_blocks[block_idx].image_types.append(image_type)

            print(f"   ✅ Captions done in {time.time() - caption_start:.2f}s")

        # Content detection and extraction for each block
        from lib.content_detector import (
            detect_content_type,
            extract_definitions,
            extract_procedure_steps,
        )
        from lib.math_to_speech import extract_equations

        for block in content_blocks:
            # Content type classification
            block.content_type = detect_content_type(block.text_content).value

            # Extract definitions (if content is definition-type or contains definition patterns)
            if block.content_type == "definition" or "definition" in block.text_content.lower():
                block.definitions = extract_definitions(block.text_content)

            # Extract procedure steps
            if block.content_type == "procedure":
                block.procedure_steps = extract_procedure_steps(block.text_content)

            # Extract equations with spoken forms
            block.equations = extract_equations(block.text_content)

            # Extract code blocks
            block.code_blocks = self._extract_code_blocks(block.text_content)

            # Extract structured tables
            for table in block.related_tables:
                table_data = self._extract_table_structured(table)
                if table_data:
                    block.tables.append(table_data)

        # Build combined context for all blocks
        for block in content_blocks:
            block.combined_context = self._combine_context(block)

        elapsed = time.time() - start_time
        print(f"✅ PDF extraction completed in {elapsed:.2f}s ({len(content_blocks)} blocks)")
        return content_blocks
    
    def _extract_scanned_pdf(self, pdf_path: str, start_time: float) -> List[ContentBlock]:
        """
        Extract scanned PDF content using OCR.
        Uses Gemini for handwriting, olmOCR for printed scans.
        """
        from unstructured.partition.text import partition_text
        
        # Try Gemini first (better for handwriting), falls back to olmOCR
        ocr_text = ocr_with_gemini(pdf_path)
        
        if not ocr_text:
            print("   ⚠️ OCR returned no text")
            return []
        
        print(f"   ✅ OCR extracted {len(ocr_text)} chars")
        
        # Chunk the OCR text using unstructured
        chunks = partition_text(
            text=ocr_text,
            chunking_strategy="by_title",
            max_characters=CHUNK_MAX_CHARS,
            combine_text_under_n_chars=CHUNK_COMBINE_UNDER,
            new_after_n_chars=CHUNK_NEW_AFTER,
            overlap=CHUNK_OVERLAP,
        )
        
        # Convert to ContentBlocks
        from lib.content_detector import (
            detect_content_type,
            extract_definitions,
            extract_procedure_steps,
        )
        from lib.math_to_speech import extract_equations
        import re

        content_blocks = []
        for idx, chunk in enumerate(chunks):
            block = ContentBlock()
            block.chunk_index = idx
            block.text_content = str(chunk)

            # Try to extract page number from text markers (=== PAGE X ===)
            page_match = re.search(r'=== PAGE (\d+) ===', block.text_content)
            if page_match:
                block.page_number = int(page_match.group(1))
            else:
                block.page_number = idx + 1

            # Content type classification
            block.content_type = detect_content_type(block.text_content).value

            # Extract definitions
            if block.content_type == "definition" or "definition" in block.text_content.lower():
                block.definitions = extract_definitions(block.text_content)

            # Extract procedure steps
            if block.content_type == "procedure":
                block.procedure_steps = extract_procedure_steps(block.text_content)

            # Extract equations with spoken forms
            block.equations = extract_equations(block.text_content)

            # Extract code blocks
            block.code_blocks = self._extract_code_blocks(block.text_content)

            block.combined_context = block.text_content
            content_blocks.append(block)

        elapsed = time.time() - start_time
        print(f"✅ Scanned PDF extraction completed in {elapsed:.2f}s ({len(content_blocks)} blocks)")
        return content_blocks
    
    def _extract_docx(self, docx_path: str) -> List[ContentBlock]:
        """
        Extract DOCX content using python-docx (lightweight, no Tesseract).

        Pipeline:
        1. python-docx for text extraction (paragraphs + tables)
        2. → Unified pipeline (chunking, hierarchy, topic blocks, enrichment)
        """
        from docx import Document

        print("⏳ Starting DOCX extraction (python-docx)...")
        start_time = time.time()

        # Step 1: Extract text with python-docx
        doc = Document(docx_path)

        # Extract paragraphs and tables
        elements = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    elements.append(f"\n## {text}\n")
                else:
                    elements.append(text)

        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(row_text)
            if table_text:
                elements.append("\n" + "\n".join(table_text) + "\n")

        # Group into ~4000 char "pages" for consistent chunking
        pages_text = []
        current_page = []
        current_length = 0

        for text in elements:
            if current_length + len(text) > 4000 and current_page:
                pages_text.append("\n\n".join(current_page))
                current_page = [text]
                current_length = len(text)
            else:
                current_page.append(text)
                current_length += len(text)

        if current_page:
            pages_text.append("\n\n".join(current_page))

        print(f"   ✅ Extracted {len(pages_text)} text segments")

        # Step 2: Unified pipeline
        return self._process_text_to_blocks(pages_text, "DOCX", start_time)

    def _extract_pptx(self, pptx_path: str) -> List[ContentBlock]:
        """
        Extract PPTX content using python-pptx (lightweight, no Tesseract).

        Pipeline:
        1. python-pptx for text extraction (slide-by-slide)
        2. → Unified pipeline (chunking, hierarchy, topic blocks, enrichment)
        """
        from pptx import Presentation

        print("⏳ Starting PPTX extraction (python-pptx)...")
        start_time = time.time()

        # Step 1: Extract text with python-pptx
        prs = Presentation(pptx_path)

        pages_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                # Extract text from text frames
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)

            # Combine slide text
            if slide_texts:
                slide_text = f"## Slide {slide_num}\n\n" + "\n\n".join(slide_texts)
                pages_text.append(slide_text)

        print(f"   ✅ Extracted {len(pages_text)} slides")

        # Step 2: Unified pipeline
        return self._process_text_to_blocks(pages_text, "PPTX", start_time)

    def _extract_markdown(self, md_path: str) -> List[ContentBlock]:
        """
        Extract Markdown content, then unified pipeline.

        Pipeline:
        1. Read markdown file
        2. → Unified pipeline (chunking, hierarchy, topic blocks, enrichment)
        """
        print("⏳ Starting Markdown extraction...")
        start_time = time.time()

        # Step 1: Read markdown file directly
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Split by major headers (##) to create "pages"
        import re
        sections = re.split(r'\n(?=#{1,2}\s)', content)
        pages_text = [s.strip() for s in sections if s.strip()]

        if not pages_text:
            pages_text = [content]

        print(f"   ✅ Extracted {len(pages_text)} sections")

        # Step 2: Unified pipeline
        return self._process_text_to_blocks(pages_text, "Markdown", start_time)

    def _extract_image(self, image_path: str) -> List[ContentBlock]:
        """
        Extract text from image using OCR, then unified pipeline.

        Pipeline:
        1. OCR extraction (auto-detects handwritten vs printed)
        2. → Unified pipeline (chunking, hierarchy, topic blocks, enrichment)

        Supports: PNG, JPEG, TIFF, BMP, WEBP
        """
        from lib.encoding_check import extract_image_with_ocr

        print("⏳ Starting image OCR extraction...")
        start_time = time.time()

        # Step 1: Extract text using OCR (auto-detects handwriting)
        pages_text = extract_image_with_ocr(image_path, ocr_provider="auto")

        if not pages_text or not pages_text[0]:
            print("⚠️ No text extracted from image")
            return []

        print(f"   ✅ OCR extracted {len(pages_text[0])} characters")

        # Step 2: Unified pipeline (even for single image, gives consistent structure)
        return self._process_text_to_blocks(pages_text, "Image", start_time)

    def _extract_auto(self, file_path: str) -> List[ContentBlock]:
        """
        Auto-detect format and extract content using Unstructured.

        NOTE: This is a fallback method. Requires Unstructured + Tesseract.
        Prefer using format-specific extractors (PDF, DOCX, PPTX, MD, Image).
        """
        # Import Unstructured locally (optional dependency)
        try:
            from unstructured.partition.auto import partition
        except ImportError:
            raise ImportError(
                "Unstructured is required for auto-detect extraction. "
                "Install with: pip install unstructured"
            )

        print("⏳ Starting auto-detect extraction (Unstructured)...")
        start_time = time.time()

        chunks = partition(
            filename=file_path,
            chunking_strategy="by_title",
            max_characters=CHUNK_MAX_CHARS,
            combine_text_under_n_chars=CHUNK_COMBINE_UNDER,
            new_after_n_chars=CHUNK_NEW_AFTER,
            overlap=CHUNK_OVERLAP,
        )

        return self._chunks_to_content_blocks(chunks, start_time, "Auto")
    
    def _chunks_to_content_blocks(self, chunks, start_time: float, format_name: str) -> List[ContentBlock]:
        """Convert unstructured chunks to ContentBlocks with rich metadata extraction."""
        from lib.content_detector import (
            detect_content_type,
            extract_definitions,
            extract_procedure_steps,
        )
        from lib.math_to_speech import extract_equations
        from lib.voice_optimizer import table_to_speech, code_to_speech

        parse_time = time.time() - start_time
        print(f"   ✅ {format_name} parsed in {parse_time:.2f}s")

        content_blocks = []
        for idx, chunk in enumerate(chunks):
            block = ContentBlock()
            block.chunk_index = idx
            block.text_content = str(chunk)

            # Extract page/slide number if available
            if hasattr(chunk, "metadata"):
                if hasattr(chunk.metadata, "page_number") and chunk.metadata.page_number:
                    block.page_number = chunk.metadata.page_number
                elif hasattr(chunk.metadata, "slide_number") and chunk.metadata.slide_number:
                    block.page_number = chunk.metadata.slide_number

            # Extract structured tables
            if "Table" in str(type(chunk)):
                block.related_tables = [chunk]
                table_data = self._extract_table_structured(chunk)
                if table_data:
                    block.tables.append(table_data)
                    block.table_descriptions.append(table_data.get("description", ""))

            # Content type classification
            block.content_type = detect_content_type(block.text_content).value

            # Extract definitions (if content is definition-type or contains definition patterns)
            if block.content_type == "definition" or "definition" in block.text_content.lower():
                block.definitions = extract_definitions(block.text_content)

            # Extract procedure steps
            if block.content_type == "procedure":
                block.procedure_steps = extract_procedure_steps(block.text_content)

            # Extract equations with spoken forms
            block.equations = extract_equations(block.text_content)

            # Extract code blocks
            block.code_blocks = self._extract_code_blocks(block.text_content)

            content_blocks.append(block)

        # Build combined context for all blocks
        for block in content_blocks:
            block.combined_context = self._combine_context(block)

        elapsed = time.time() - start_time
        print(f"✅ {format_name} extraction completed in {elapsed:.2f}s ({len(content_blocks)} blocks)")
        return content_blocks

    def _extract_table_structured(self, table_element) -> Optional[Dict]:
        """Extract table as structured data (headers, rows, spoken form)."""
        from lib.voice_optimizer import table_to_speech

        try:
            rows = []
            headers = []

            if hasattr(table_element, 'metadata') and hasattr(table_element.metadata, 'text_as_html'):
                html = table_element.metadata.text_as_html
                if html:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(html, 'html.parser')

                    # Get headers from first row
                    header_row = soup.find('tr')
                    if header_row:
                        headers = [th.get_text(strip=True) for th in header_row.find_all(['th', 'td'])]

                    # Get data rows
                    for tr in soup.find_all('tr')[1:]:
                        row = [td.get_text(strip=True) for td in tr.find_all('td')]
                        if row:
                            rows.append(row)

            # Fallback: try to parse from text
            if not headers and not rows:
                text = str(table_element)
                lines = [l.strip() for l in text.split('\n') if l.strip()]
                if len(lines) >= 2:
                    headers = lines[0].split('\t') if '\t' in lines[0] else lines[0].split('|')
                    headers = [h.strip() for h in headers if h.strip()]
                    for line in lines[1:]:
                        cells = line.split('\t') if '\t' in line else line.split('|')
                        cells = [c.strip() for c in cells if c.strip()]
                        if cells:
                            rows.append(cells)

            if headers or rows:
                spoken = table_to_speech(headers, rows) if headers else ""
                return {
                    "headers": headers,
                    "rows": rows,
                    "row_count": len(rows),
                    "col_count": len(headers),
                    "description": f"Table with {len(headers)} columns and {len(rows)} rows",
                    "spoken": spoken,
                }
        except Exception as e:
            print(f"   ⚠️ Table extraction error: {e}")

        return None

    def _extract_code_blocks(self, text: str) -> List[Dict]:
        """Extract code blocks with language detection and descriptions."""
        from lib.voice_optimizer import code_to_speech
        import re

        code_blocks = []

        # Markdown code blocks: ```language\ncode\n```
        pattern = r'```(\w*)\n(.*?)```'
        matches = re.findall(pattern, text, re.DOTALL)

        for language, code in matches:
            if not language:
                language = self._detect_code_language(code)

            code = code.strip()
            if code:
                code_blocks.append({
                    "language": language or "unknown",
                    "code": code,
                    "line_count": len(code.split('\n')),
                    "description": f"{language.capitalize() if language else 'Code'} block with {len(code.split(chr(10)))} lines",
                    "spoken": code_to_speech(code, language or "unknown"),
                })

        return code_blocks

    def _detect_code_language(self, code: str) -> str:
        """Detect programming language from code heuristics."""
        code_lower = code.lower()

        # Python indicators
        if 'def ' in code or 'import ' in code or 'from ' in code or '__name__' in code:
            return 'python'
        # JavaScript indicators
        if 'function ' in code or 'const ' in code or 'let ' in code or '=>' in code:
            return 'javascript'
        # Java/C++ indicators
        if 'public ' in code or 'private ' in code or 'class ' in code and 'def' not in code:
            return 'java'
        # SQL indicators
        if 'SELECT ' in code.upper() or 'FROM ' in code.upper() or 'WHERE ' in code.upper():
            return 'sql'
        # Shell indicators
        if code.startswith('#!') or 'echo ' in code_lower or 'export ' in code_lower:
            return 'bash'

        return 'unknown'
    
    # Backward compatibility alias
    def extract_pdf(self, pdf_path: str) -> List[ContentBlock]:
        """Backward compatible: Extract PDF (use extract_document for all formats)"""
        return self._extract_pdf(pdf_path)
    
    def enrich_content_blocks(
        self,
        content_blocks: List[ContentBlock],
        parallel_questions: bool = True,
        max_workers: int = 4
    ) -> List[ContentBlock]:
        """
        Enrich content blocks with embeddings and questions.
        Uses parallel processing for significant speedup.

        Args:
            content_blocks: List of ContentBlocks to enrich
            parallel_questions: Whether to generate questions in parallel
            max_workers: Max concurrent question generation workers
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed

        total_start = time.time()
        print(f"\n⏳ Enriching {len(content_blocks)} content blocks...")

        # Step 1: Build combined context for all blocks (fast, sequential)
        for block in content_blocks:
            block.combined_context = self._combine_context(block)

        # Step 2: Generate embeddings in parallel (fast, ~1s each)
        print(f"🔢 Generating embeddings for {len(content_blocks)} blocks in parallel...")
        embed_start = time.time()

        def generate_embedding(block_idx: int) -> tuple:
            block = content_blocks[block_idx]
            if block.combined_context:
                embedding = embed_text(block.combined_context)
                return block_idx, embedding
            return block_idx, []

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [executor.submit(generate_embedding, i) for i in range(len(content_blocks))]
            for future in as_completed(futures):
                idx, embedding = future.result()
                content_blocks[idx].embeddings = embedding

        print(f"   ✅ All embeddings done in {time.time() - embed_start:.2f}s")

        # Step 3: Generate questions (slow, ~30-40s each - parallelize!)
        if parallel_questions:
            print(f"❓ Generating questions for {len(content_blocks)} blocks in parallel (max {max_workers} workers)...")
            question_start = time.time()

            def generate_questions_for_block(block_idx: int) -> tuple:
                block = content_blocks[block_idx]
                start = time.time()
                questions = self._generate_questions(block)
                elapsed = time.time() - start
                print(f"   ✅ Block {block_idx + 1}: {len(questions)} questions in {elapsed:.1f}s")
                return block_idx, questions

            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = [executor.submit(generate_questions_for_block, i) for i in range(len(content_blocks))]
                for future in as_completed(futures):
                    idx, questions = future.result()
                    content_blocks[idx].questions = questions

            print(f"   ✅ All questions done in {time.time() - question_start:.2f}s")
        else:
            # Sequential fallback
            for i, block in enumerate(content_blocks):
                question_start = time.time()
                print(f"❓ Generating questions for block {i+1}/{len(content_blocks)}...")
                block.questions = self._generate_questions(block)
                print(f"   ✅ Questions done in {time.time() - question_start:.2f}s ({len(block.questions)} questions)")

        total_elapsed = time.time() - total_start
        total_questions = sum(len(b.questions) for b in content_blocks)
        print(f"✅ All {len(content_blocks)} blocks enriched in {total_elapsed:.2f}s ({total_questions} questions)")
        return content_blocks

    def _combine_context(self, block: ContentBlock) -> str:
        """Combine text, image captions, and table descriptions into a single context string"""
        components = [block.text_content]
        if block.image_captions:
            components.append("\n".join(block.image_captions))
        if block.table_descriptions:
            components.append("\n".join(block.table_descriptions))
        return "\n\n".join(components)
    def _generate_questions(self, block: ContentBlock, max_retries: int = 2) -> List[Question]:
        """
        Generate structured questions for a topic-level content block.

        Generates exactly 12 questions per topic:
        - 2 questions per Bloom's taxonomy level (Remember, Understand, Apply, Analyze, Evaluate, Create)
        - Mix of long-answer (4) and multiple-choice (8) questions
        - Balanced difficulty distribution

        Includes:
        - Retry logic for API failures
        - JSON parsing with error recovery
        - Figure reference validation
        - Graceful handling of invalid enum values
        - TTS-optimized versions (spoken_text, spoken_options)
        """
        import httpx
        import json
        from lib.voice_optimizer import optimize_for_tts

        # Build figure info for prompt if figures exist
        has_figures = hasattr(block, 'figure_map') and block.figure_map
        max_figure = max(block.figure_map.keys()) if has_figures else 0
        figure_instruction = ""
        if has_figures:
            figure_instruction = f"""
IMPORTANT - This content includes {len(block.figure_map)} figure(s):
- When creating questions about figures, reference them as "Figure 1", "Figure 2", etc.
- Include "figure_ref": N in the question JSON to link it to the correct figure
- Valid figure numbers are 1 to {max_figure}
- At least 1-2 questions should test understanding of the figures if they contain educational content
"""

        # Topic info for context
        topic_info = ""
        if block.section_title:
            topic_info = f"\nTOPIC: {block.section_title}"
            if block.chapter_title:
                topic_info += f" (from {block.chapter_title})"

        prompt = f"""You are an expert educational assessment designer. Generate balanced questions across Bloom's taxonomy.
{topic_info}
CONTENT:
{block.combined_context[:8000]}
{figure_instruction}
REQUIREMENTS - Generate exactly 12 questions (2 per Bloom's level):

BLOOM'S TAXONOMY - Generate EXACTLY 2 questions for EACH level:
1. REMEMBER (2 questions) - Recall facts, definitions, terms
   - Question types: "What is...?", "Define...", "List...", "Name..."
   - Difficulty: basic

2. UNDERSTAND (2 questions) - Explain ideas, summarize, interpret
   - Question types: "Explain...", "Describe...", "Summarize...", "Compare..."
   - Difficulty: basic to intermediate

3. APPLY (2 questions) - Use knowledge in new situations, solve problems
   - Question types: "Calculate...", "Solve...", "Demonstrate...", "How would you use..."
   - Difficulty: intermediate

4. ANALYZE (2 questions) - Examine relationships, draw connections
   - Question types: "Why does...?", "What is the relationship...?", "Differentiate..."
   - Difficulty: intermediate to advanced

5. EVALUATE (2 questions) - Judge, critique, assess validity
   - Question types: "Assess...", "Critique...", "Justify...", "Which is better..."
   - Difficulty: advanced

6. CREATE (2 questions) - Design, construct, propose new solutions
   - Question types: "Design...", "Propose...", "What would happen if...", "Construct..."
   - Difficulty: advanced

QUESTION TYPE DISTRIBUTION:
- 4 Long-answer questions (for Analyze, Evaluate, Create levels)
- 8 Multiple-choice questions (for Remember, Understand, Apply levels)

MATH/LATEX HANDLING (CRITICAL):
- The content contains LaTeX math notation like $\\Delta Q$, $10^{{22}}$, $$E = mc^2$$
- PRESERVE LaTeX format EXACTLY in question text and options
- Use $...$ for inline math, $$...$$ for display equations
- Do NOT convert LaTeX to plain text (wrong: "delta Q", correct: "$\\Delta Q$")
- Do NOT use Unicode approximations (wrong: "ΔQ", correct: "$\\Delta Q$")
- Example question: "What is the change in charge $\\Delta Q$ when..."
- Example option: "A. $v = \\frac{{d}}{{t}}$"

FOR ALL QUESTIONS:
- text: Clear, unambiguous question text (preserve LaTeX math notation)
- bloom_level: remember/understand/apply/analyze/evaluate/create
- difficulty: basic/intermediate/advanced
- key_points: 2-4 points the answer should cover (CRITICAL FOR GRADING)

FOR LONG ANSWER:
- question_type: "long_answer"
- expected_time: 5-10 minutes

FOR MULTIPLE CHOICE:
- question_type: "multiple_choice"
- expected_time: 1-3 minutes
- options: Exactly 4 options as full text
- correct_answer: The letter (A, B, C, or D)
- DISTRACTOR REQUIREMENTS:
  * Each wrong option represents a plausible misconception
  * Base on common student errors
  * All options similar in length and structure

Output ONLY valid JSON:
{{
  "questions": [
    {{"text": "...", "bloom_level": "remember", "difficulty": "basic", "question_type": "multiple_choice", "expected_time": 2, "key_points": ["..."], "options": ["A...", "B...", "C...", "D..."], "correct_answer": "A", "figure_ref": null}},
    {{"text": "...", "bloom_level": "remember", "difficulty": "basic", "question_type": "multiple_choice", "expected_time": 2, "key_points": ["..."], "options": ["A...", "B...", "C...", "D..."], "correct_answer": "B", "figure_ref": null}},
    {{"text": "...", "bloom_level": "understand", "difficulty": "intermediate", "question_type": "multiple_choice", "expected_time": 2, "key_points": ["..."], "options": ["..."], "correct_answer": "C", "figure_ref": null}},
    {{"text": "...", "bloom_level": "understand", "difficulty": "intermediate", "question_type": "multiple_choice", "expected_time": 3, "key_points": ["..."], "options": ["..."], "correct_answer": "A", "figure_ref": null}},
    {{"text": "...", "bloom_level": "apply", "difficulty": "intermediate", "question_type": "multiple_choice", "expected_time": 3, "key_points": ["..."], "options": ["..."], "correct_answer": "B", "figure_ref": null}},
    {{"text": "...", "bloom_level": "apply", "difficulty": "intermediate", "question_type": "multiple_choice", "expected_time": 3, "key_points": ["..."], "options": ["..."], "correct_answer": "D", "figure_ref": null}},
    {{"text": "...", "bloom_level": "analyze", "difficulty": "advanced", "question_type": "multiple_choice", "expected_time": 3, "key_points": ["..."], "options": ["..."], "correct_answer": "A", "figure_ref": null}},
    {{"text": "...", "bloom_level": "analyze", "difficulty": "advanced", "question_type": "long_answer", "expected_time": 7, "key_points": ["..."], "figure_ref": null}},
    {{"text": "...", "bloom_level": "evaluate", "difficulty": "advanced", "question_type": "long_answer", "expected_time": 8, "key_points": ["..."], "figure_ref": null}},
    {{"text": "...", "bloom_level": "evaluate", "difficulty": "advanced", "question_type": "long_answer", "expected_time": 7, "key_points": ["..."], "figure_ref": null}},
    {{"text": "...", "bloom_level": "create", "difficulty": "advanced", "question_type": "long_answer", "expected_time": 10, "key_points": ["..."], "figure_ref": null}},
    {{"text": "...", "bloom_level": "create", "difficulty": "advanced", "question_type": "multiple_choice", "expected_time": 3, "key_points": ["..."], "options": ["..."], "correct_answer": "C", "figure_ref": null}}
  ]
}}

Note: "figure_ref" should be the figure number (1 to {max_figure if max_figure else 'N'}) if the question references a figure, or null if not.
"""

        def parse_bloom_level(value: str) -> BloomLevel:
            """Safely parse bloom level with fallback."""
            try:
                return BloomLevel(value.lower())
            except (ValueError, AttributeError):
                return BloomLevel.UNDERSTAND  # Safe default

        def parse_difficulty(value: str) -> Difficulty:
            """Safely parse difficulty with fallback."""
            try:
                return Difficulty(value.lower())
            except (ValueError, AttributeError):
                return Difficulty.INTERMEDIATE  # Safe default

        def resolve_figure_ref(figure_ref, block, has_figures, max_figure) -> tuple:
            """Resolve figure_ref to URL and description, with validation."""
            if not figure_ref or not has_figures:
                return None, None, None

            # Validate figure_ref is an integer
            if not isinstance(figure_ref, int):
                try:
                    figure_ref = int(figure_ref)
                except (ValueError, TypeError):
                    return None, None, None

            # Validate figure_ref is in bounds
            if figure_ref < 1 or figure_ref > max_figure:
                # Out of bounds - clear the reference
                return None, None, None

            if figure_ref in block.figure_map:
                fig_data = block.figure_map[figure_ref]
                return figure_ref, fig_data.get("url"), fig_data.get("description")

            return None, None, None

        def parse_json_response(text: str) -> dict:
            """Parse JSON from LLM response with multiple strategies."""
            # Strategy 1: Find JSON object directly
            start = text.find('{')
            end = text.rfind('}') + 1
            if start != -1 and end > start:
                try:
                    return json.loads(text[start:end])
                except json.JSONDecodeError:
                    pass

            # Strategy 2: Try to extract from markdown code block
            if "```json" in text:
                try:
                    json_str = text.split("```json")[1].split("```")[0]
                    return json.loads(json_str.strip())
                except (IndexError, json.JSONDecodeError):
                    pass

            if "```" in text:
                try:
                    json_str = text.split("```")[1].split("```")[0]
                    return json.loads(json_str.strip())
                except (IndexError, json.JSONDecodeError):
                    pass

            raise ValueError("Could not parse JSON from response")

        # Retry loop
        last_error = None
        for attempt in range(max_retries + 1):
            try:
                # Use Cerebras for question generation (higher rate limits with pay-as-you-go)
                api_url = CEREBRAS_URL
                api_key = CEREBRAS_API_KEY
                model = "gpt-oss-120b"

                response = httpx.post(
                    api_url,
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": model,
                        "messages": [
                            {"role": "system", "content": "You are an educational question generator. Output valid JSON only."},
                            {"role": "user", "content": prompt}
                        ],
                        "max_tokens": 8000,
                        "temperature": 0.3
                    },
                    timeout=60.0
                )

                if response.status_code == 429:
                    # Rate limited - wait and retry
                    if attempt < max_retries:
                        time.sleep(2 ** attempt)
                        continue
                    print(f"   ⚠️ Rate limited after {max_retries + 1} attempts")
                    return []

                if response.status_code != 200:
                    last_error = f"API error: {response.status_code}"
                    if attempt < max_retries:
                        time.sleep(1)
                        continue
                    print(f"   ⚠️ GPT OSS error: {response.status_code}")
                    return []

                result = response.json()
                msg = result["choices"][0]["message"]
                # gpt-oss-120b returns 'reasoning' instead of 'content'
                text = msg.get("content") or msg.get("reasoning", "")

                # Parse JSON with multiple strategies
                data = parse_json_response(text)

                # Convert to Question objects
                questions = []

                # Handle new unified "questions" format (12 questions per topic)
                all_questions = data.get("questions", [])

                # Fallback: support old format with separate arrays
                if not all_questions:
                    all_questions = data.get("long_answer_questions", []) + data.get("multiple_choice_questions", [])

                for q in all_questions:
                    figure_ref, image_url, image_desc = resolve_figure_ref(
                        q.get("figure_ref"), block, has_figures, max_figure
                    )

                    question_text = q.get("text", "")
                    question_type_str = q.get("question_type", "multiple_choice").lower()
                    options = q.get("options", [])

                    # Determine question type
                    if question_type_str == "long_answer":
                        q_type = QuestionType.LONG_ANSWER
                    else:
                        q_type = QuestionType.MULTIPLE_CHOICE

                    try:
                        questions.append(Question(
                            text=question_text,
                            bloom_level=parse_bloom_level(q.get("bloom_level", "understand")),
                            difficulty=parse_difficulty(q.get("difficulty", "intermediate")),
                            question_type=q_type,
                            expected_time=q.get("expected_time", 5 if q_type == QuestionType.LONG_ANSWER else 2),
                            key_points=q.get("key_points", []),
                            options=options if q_type == QuestionType.MULTIPLE_CHOICE else None,
                            correct_answer=q.get("correct_answer") if q_type == QuestionType.MULTIPLE_CHOICE else None,
                            # TTS: questions are read verbatim
                            spoken_text=optimize_for_tts(question_text),
                            spoken_options=[optimize_for_tts(opt) for opt in options] if options and q_type == QuestionType.MULTIPLE_CHOICE else None,
                            figure_ref=figure_ref,
                            image_url=image_url,
                            image_description=image_desc
                        ))
                    except Exception as e:
                        print(f"   ⚠️ Skipping malformed question: {e}")

                # Log Bloom's distribution for verification
                bloom_counts = {}
                for q in questions:
                    level = q.bloom_level.value
                    bloom_counts[level] = bloom_counts.get(level, 0) + 1
                print(f"   📊 Generated {len(questions)} questions: {bloom_counts}")

                return questions

            except json.JSONDecodeError as e:
                last_error = f"JSON parse error: {e}"
                if attempt < max_retries:
                    time.sleep(1)
                    continue
            except Exception as e:
                last_error = str(e)
                if attempt < max_retries:
                    time.sleep(1)
                    continue

        print(f"   ⚠️ Failed to generate questions after {max_retries + 1} attempts: {last_error}")
        return []

    def _extract_image_base64_from_chunk(self, chunk) -> List[str]:
        """Extract base64 image data from chunk without captioning"""
        images = []
        if hasattr(chunk, 'metadata') and hasattr(chunk.metadata, 'orig_elements'):
            for el in chunk.metadata.orig_elements:
                if "Image" in str(type(el)):
                    if hasattr(el.metadata, 'image_base64') and el.metadata.image_base64:
                        try:
                            import base64
                            base64.b64decode(el.metadata.image_base64)  # Validate
                            images.append(el.metadata.image_base64)
                        except Exception as e:
                            print(f"Invalid image base64: {e}")
                            continue
        return images
    
    def _batch_caption_images(self, images: List[str], max_workers: int = 5) -> List[Tuple[str, str]]:
        """
        Caption multiple images in parallel and classify their types.
        GPT-4o-mini handles ~5 concurrent requests well within rate limits.

        Returns:
            List of tuples: (caption, image_type)
            Image types: circuit, anatomy, flowchart, graph, diagram, equation, map, table, photo, other
        """
        def caption_and_classify(img_b64: str) -> Tuple[str, str]:
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": """Analyze this educational image and provide:
1. A detailed description focusing on scientific concepts, diagrams, data, or educational content.
2. Classify the image type as ONE of: circuit, anatomy, flowchart, graph, diagram, equation, map, table, photo, other

Format your response as:
TYPE: [one of the types above]
DESCRIPTION: [your detailed description]"""
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}
                        }
                    ]
                }
            ]
            try:
                response = self.vision_llm.invoke(messages)
                content = response.content

                # Parse response to extract type and description
                image_type = "other"
                description = content

                if "TYPE:" in content:
                    lines = content.split("\n")
                    for line in lines:
                        if line.strip().startswith("TYPE:"):
                            type_value = line.split(":", 1)[1].strip().lower()
                            # Validate type
                            valid_types = ["circuit", "anatomy", "flowchart", "graph", "diagram", "equation", "map", "table", "photo", "other"]
                            if type_value in valid_types:
                                image_type = type_value
                        elif line.strip().startswith("DESCRIPTION:"):
                            description = line.split(":", 1)[1].strip()
                            # Get remaining lines as part of description
                            idx = lines.index(line)
                            if idx < len(lines) - 1:
                                description += " " + " ".join(lines[idx+1:])
                            break

                return (description, image_type)
            except Exception as e:
                print(f"Failed to caption image: {e}")
                return ("Image description unavailable", "other")

        # Parallel execution
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            results = list(executor.map(caption_and_classify, images))

        return results
    
    def _extract_tables_from_chunk(self, chunk) -> List:
        """Extract table objects from a CompositeElement chunk"""
        tables = []
        if hasattr(chunk, 'metadata') and hasattr(chunk.metadata, 'orig_elements'):
            for el in chunk.metadata.orig_elements:
                if "Table" in str(type(el)):
                    tables.append(el)
        return tables
    
    # ========== Neo4j Write Functions ==========
    
    def _upsert_document(self, tx, doc_meta: dict):
        """Create or update document node in Neo4j"""
        tx.run(
            """
            MERGE (u:User {id: $user_id})
            MERGE (d:Document {documentId: $doc_id})
            ON CREATE SET 
                d.title = $title,
                d.source = $source,
                d.created_at = datetime()
            MERGE (u)-[:UPLOADED]->(d)
            """,
            **doc_meta
        )
    
    def _create_content_block(self, tx, payload: dict):
        """Create ContentBlock node with embeddings, hierarchy, and rich metadata"""
        tx.run(
            """
            MATCH (d:Document {documentId: $doc_id})
            CREATE (cb:ContentBlock {
                block_id: $block_id,
                chunk_index: $chunk_index,
                text_content: $text_content,
                combined_context: $combined_context,
                page_from: $page_from,
                page_to: $page_to,
                page_start: $page_start,
                page_end: $page_end,
                has_images: $has_images,
                has_tables: $has_tables,
                image_count: $image_count,
                table_count: $table_count,
                embedding: $embedding,
                page_number: $page_number,
                bbox: $bbox,
                chapter_title: $chapter_title,
                section_title: $section_title,
                heading_level: $heading_level,
                image_urls: $image_urls,
                image_descriptions: $image_descriptions,
                image_types: $image_types,
                content_type: $content_type,
                definitions: $definitions,
                procedure_steps: $procedure_steps,
                equations: $equations,
                code_blocks: $code_blocks,
                tables: $tables
            })
            MERGE (d)-[:HAS_CONTENT_BLOCK]->(cb)
            """,
            **payload
        )

    def _create_hierarchy_nodes(self, tx, doc_id: str, content_blocks: List[ContentBlock]):
        """
        Create Chapter and Section nodes in Neo4j for hierarchy-aware retrieval.
        Links ContentBlocks to their respective Chapter/Section.
        """
        # Collect unique chapters and sections
        chapters = {}  # {chapter_title: [block_ids]}
        sections = {}  # {(chapter_title, section_title): [block_ids]}

        for i, block in enumerate(content_blocks):
            block_id = f"{doc_id}::block::{i}"

            if block.chapter_title:
                if block.chapter_title not in chapters:
                    chapters[block.chapter_title] = []
                chapters[block.chapter_title].append(block_id)

                if block.section_title:
                    key = (block.chapter_title, block.section_title)
                    if key not in sections:
                        sections[key] = []
                    sections[key].append(block_id)

        # Create Chapter nodes
        for chapter_title, block_ids in chapters.items():
            tx.run(
                """
                MATCH (d:Document {documentId: $doc_id})
                MERGE (c:Chapter {doc_id: $doc_id, title: $chapter_title})
                MERGE (d)-[:HAS_CHAPTER]->(c)
                WITH c
                UNWIND $block_ids AS block_id
                MATCH (cb:ContentBlock {block_id: block_id})
                MERGE (c)-[:CONTAINS]->(cb)
                """,
                doc_id=doc_id,
                chapter_title=chapter_title,
                block_ids=block_ids
            )

        # Create Section nodes
        for (chapter_title, section_title), block_ids in sections.items():
            tx.run(
                """
                MATCH (c:Chapter {doc_id: $doc_id, title: $chapter_title})
                MERGE (s:Section {doc_id: $doc_id, title: $section_title, chapter: $chapter_title})
                MERGE (c)-[:HAS_SECTION]->(s)
                WITH s
                UNWIND $block_ids AS block_id
                MATCH (cb:ContentBlock {block_id: block_id})
                MERGE (s)-[:CONTAINS]->(cb)
                """,
                doc_id=doc_id,
                chapter_title=chapter_title,
                section_title=section_title,
                block_ids=block_ids
            )

    def _create_question_set(self, tx, block_id: str, questions: List[Question], doc_id: str):
        """
        Create a single QuestionSet node containing all questions for a ContentBlock.
        Stores questions as JSON array instead of separate nodes (78% node reduction).
        """
        from datetime import datetime
        
        # Convert questions to JSON-serializable format
        questions_json = json.dumps([
            {
                "question_id": f"{block_id}::q::{idx}",
                "text": q.text,
                "bloom_level": q.bloom_level.value,
                "difficulty": q.difficulty.value,
                "question_type": q.question_type.value,
                "expected_time": q.expected_time,
                "key_points": q.key_points or [],
                "options": q.options or [],
                "correct_answer": q.correct_answer or "",
                "image_url": q.image_url,
                "image_description": q.image_description
            }
            for idx, q in enumerate(questions)
        ])
        
        # Calculate difficulty distribution
        difficulty_dist = {}
        for q in questions:
            diff = q.difficulty.value
            difficulty_dist[diff] = difficulty_dist.get(diff, 0) + 1

        # Calculate Bloom's taxonomy distribution
        bloom_dist = {}
        for q in questions:
            bloom = q.bloom_level.value
            bloom_dist[bloom] = bloom_dist.get(bloom, 0) + 1

        tx.run(
            """
            MATCH (cb:ContentBlock {block_id: $block_id})
            CREATE (qs:QuestionSet {
                questionset_id: $questionset_id,
                questions: $questions_json,
                total_count: $total_count,
                difficulty_distribution: $difficulty_dist,
                bloom_distribution: $bloom_dist,
                generated_at: $generated_at,
                doc_id: $doc_id
            })
            CREATE (cb)-[:HAS_QUESTIONS]->(qs)
            """,
            block_id=block_id,
            questionset_id=f"{block_id}::qs",
            questions_json=questions_json,
            total_count=len(questions),
            difficulty_dist=json.dumps(difficulty_dist),
            bloom_dist=json.dumps(bloom_dist),
            generated_at=datetime.utcnow().isoformat(),
            doc_id=doc_id
        )
    
    def persist_to_neo4j(self, doc_id: str, doc_meta: dict, content_blocks: List[ContentBlock]):
        """
        Persist document, content blocks, and questions to Neo4j.
        Creates the full graph structure with embeddings.
        """
        if not self.neo4j_driver:
            print("⚠️  Neo4j not available, skipping graph persistence")
            return
        
        persist_start = time.time()
        print("\n💾 Persisting to Neo4j...")
        
        with self.neo4j_driver.session() as session:
            # 1. Create/update document
            session.execute_write(self._upsert_document, doc_meta)
            print(f"✅ Created document: {doc_id}")
            
            # 2. Create content blocks
            block_ids = []
            skipped_blocks = 0
            for i, block in enumerate(content_blocks):
                # Validation: Skip blocks with empty text_content to prevent ghost nodes
                if not block.text_content or not block.text_content.strip():
                    print(f"⚠️  Skipping block {i} - empty text_content")
                    skipped_blocks += 1
                    continue

                block_id = f"{doc_id}::block::{i}"
                block_ids.append(block_id)

                # Extract page metadata if available
                page_from = block.meta.get("page_from")
                page_to = block.meta.get("page_to")
                
                # Use page_start/page_end if available (topic-level blocks), fallback to page_number
                block_page_start = getattr(block, 'page_start', block.page_number)
                block_page_end = getattr(block, 'page_end', block.page_number)

                payload = {
                    "doc_id": doc_id,
                    "block_id": block_id,
                    "chunk_index": i,
                    "text_content": block.text_content[:5000],  # Truncate long text
                    "combined_context": block.combined_context[:5000],
                    "page_from": page_from,
                    "page_to": page_to,
                    "page_start": block_page_start,  # First page of topic
                    "page_end": block_page_end,      # Last page of topic
                    "has_images": len(block.image_urls) > 0 or len(block.image_captions) > 0,
                    "has_tables": len(block.related_tables) > 0,
                    "image_count": len(block.image_urls) or len(block.image_captions),
                    "table_count": len(block.related_tables),
                    "embedding": block.embeddings or [],
                    "page_number": block.page_number,
                    "bbox": block.bbox,  # Optional: [x1, y1, x2, y2] for scroll-to
                    # Hierarchy fields
                    "chapter_title": block.chapter_title,
                    "section_title": block.section_title,
                    "heading_level": block.heading_level,
                    # Image fields
                    "image_urls": block.image_urls,
                    "image_descriptions": block.image_descriptions,
                    "image_types": getattr(block, 'image_types', []),
                    # NEW: Content type classification
                    "content_type": getattr(block, 'content_type', 'narrative'),
                    # NEW: Extracted structured content (as JSON strings for Neo4j)
                    "definitions": json.dumps(getattr(block, 'definitions', [])),
                    "procedure_steps": json.dumps(getattr(block, 'procedure_steps', [])),
                    "equations": json.dumps(getattr(block, 'equations', [])),
                    "code_blocks": json.dumps(getattr(block, 'code_blocks', [])),
                    "tables": json.dumps(getattr(block, 'tables', [])),
                }
                
                session.execute_write(self._create_content_block, payload)
                print(f"✅ Created ContentBlock {i+1}/{len(content_blocks)}")
                
                # 3. Create QuestionSet for this block (all questions in one node)
                if block.questions:
                    session.execute_write(
                        self._create_question_set,
                        block_id,
                        block.questions,
                        doc_id
                    )
                    print(f"  ✅ Created QuestionSet with {len(block.questions)} questions")

            # 4. Create hierarchy nodes (Chapter, Section) for hierarchy-aware retrieval
            session.execute_write(self._create_hierarchy_nodes, doc_id, content_blocks)
            print("✅ Created hierarchy nodes (Chapters, Sections)")

        persist_elapsed = time.time() - persist_start
        if skipped_blocks > 0:
            print(f"⚠️  Skipped {skipped_blocks} blocks with empty text_content")
        print(f"✅ Successfully persisted to Neo4j in {persist_elapsed:.2f}s")

    def ingest_document(
        self,
        file_path: str,
        doc_id: str,
        user_id: str,
        title: str = None,
        extract_images: bool = True,
        create_hierarchy: bool = True,
        generate_questions: bool = True
    ) -> dict:
        """
        Full enhanced ingestion pipeline.

        Phases:
        1. Extract text content (Unstructured for native, OCR for scanned)
        2. Extract images with PyMuPDF (for PDFs)
        3. Upload images to R2
        4. Create hierarchy with LLM
        5. Match images to content chunks
        6. Generate embeddings and questions
        7. Link images to questions
        8. Persist to Neo4j

        Args:
            file_path: Path to document file
            doc_id: Unique document identifier
            user_id: User who uploaded the document
            title: Optional document title
            extract_images: Whether to extract and upload images
            create_hierarchy: Whether to create chapter/section hierarchy
            generate_questions: Whether to generate questions

        Returns:
            Summary dict with counts and timing
        """
        import os
        from pathlib import Path

        total_start = time.time()
        print(f"\n{'='*60}")
        print(f"🚀 Enhanced Ingestion Pipeline")
        print(f"📄 Document: {Path(file_path).name}")
        print(f"🆔 Doc ID: {doc_id}")
        print(f"{'='*60}\n")

        # Default title from filename
        if not title:
            title = Path(file_path).stem

        # ===== Phase 1: Extract text content =====
        print("⏳ Phase 1: Text Extraction...")
        content_blocks = self.extract_document(file_path)
        print(f"✅ Phase 1 complete: {len(content_blocks)} content blocks\n")

        # ===== Phase 2: Extract images (PDFs only) =====
        images_by_page = {}
        image_index = {}
        ext = os.path.splitext(file_path)[1].lower()

        if extract_images and ext == '.pdf':
            print("⏳ Phase 2: Image Extraction...")
            images_by_page = extract_images_from_pdf(file_path)

            if images_by_page:
                # Filter out non-educational images using vision LLM
                total_before = sum(len(imgs) for imgs in images_by_page.values())
                images_by_page = filter_and_describe_images(
                    images_by_page,
                    vision_model="gemini",
                    max_concurrent=5
                )
                total_after = sum(len(imgs) for imgs in images_by_page.values())
                print(f"   🎯 Vision filter: {total_after} educational, {total_before - total_after} decorative removed")

                # Upload filtered images to R2
                if images_by_page:
                    image_index = upload_images_to_r2(images_by_page, doc_id)
                    print(f"✅ Phase 2 complete: {len(image_index)} images uploaded\n")
                else:
                    print("✅ Phase 2 complete: All images filtered as decorative\n")
            else:
                print("✅ Phase 2 complete: No images found\n")
        else:
            print("⏳ Phase 2: Skipped (not a PDF or images disabled)\n")

        # ===== Phase 3: Create hierarchy =====
        if create_hierarchy and len(content_blocks) > 1:
            print("⏳ Phase 3: Hierarchy Creation...")
            hierarchy = create_hierarchy_with_llm(content_blocks)
            content_blocks = apply_hierarchy_to_chunks(content_blocks, hierarchy)
            print(f"✅ Phase 3 complete: Hierarchy applied\n")
        else:
            print("⏳ Phase 3: Skipped (single block or disabled)\n")

        # ===== Phase 4: Match images to chunks =====
        if images_by_page and image_index:
            print("⏳ Phase 4: Image-Chunk Matching...")
            content_blocks = match_images_to_chunks(
                content_blocks,
                images_by_page,
                image_index,
                vision_llm=self.vision_llm  # For descriptions
            )
            print(f"✅ Phase 4 complete\n")
        else:
            print("⏳ Phase 4: Skipped (no images)\n")

        # ===== Phase 5: Enrich (embeddings + questions) =====
        if generate_questions:
            print("⏳ Phase 5: Enrichment (embeddings + questions)...")
            content_blocks = self.enrich_content_blocks(content_blocks)

            # Link images to questions
            for block in content_blocks:
                if block.image_urls:
                    add_image_context_to_questions(block)

            print(f"✅ Phase 5 complete\n")
        else:
            # Still generate embeddings even if no questions
            print("⏳ Phase 5: Generating embeddings only...")
            for block in content_blocks:
                block.combined_context = self._combine_context(block)
                if block.combined_context:
                    block.embeddings = embed_text(block.combined_context)
            print(f"✅ Phase 5 complete\n")

        # ===== Phase 6: Persist to Neo4j =====
        print("⏳ Phase 6: Neo4j Persistence...")
        doc_meta = {
            "user_id": user_id,
            "doc_id": doc_id,
            "title": title,
            "source": Path(file_path).name
        }
        self.persist_to_neo4j(doc_id, doc_meta, content_blocks)
        print(f"✅ Phase 6 complete\n")

        # ===== Summary =====
        total_elapsed = time.time() - total_start
        total_questions = sum(len(b.questions) for b in content_blocks)
        total_images = sum(len(b.image_urls) for b in content_blocks)

        # Count chapters and sections
        chapters = set(b.chapter_title for b in content_blocks if b.chapter_title)
        sections = set(b.section_title for b in content_blocks if b.section_title)

        # Calculate page count from content blocks
        page_count = max((b.page_number for b in content_blocks if b.page_number), default=1)

        summary = {
            "doc_id": doc_id,
            "title": title,
            "content_blocks": len(content_blocks),
            "page_count": page_count,
            "chapters": len(chapters),
            "sections": len(sections),
            "questions": total_questions,
            "images_extracted": sum(len(imgs) for imgs in images_by_page.values()) if images_by_page else 0,
            "images_uploaded": len(image_index),
            "images_matched": total_images,
            "elapsed_seconds": round(total_elapsed, 2)
        }

        print(f"{'='*60}")
        print(f"✅ Ingestion Complete!")
        print(f"   📦 Blocks: {summary['content_blocks']}")
        print(f"   📖 Chapters: {summary['chapters']}, Sections: {summary['sections']}")
        print(f"   ❓ Questions: {summary['questions']}")
        print(f"   🖼️  Images: {summary['images_extracted']} extracted, {summary['images_matched']} matched")
        print(f"   ⏱️  Time: {total_elapsed:.2f}s ({total_elapsed/60:.1f} min)")
        print(f"{'='*60}\n")

        return summary

    def ingest_document_fast(
        self,
        file_path: str,
        doc_id: str,
        user_id: str,
        title: str = None,
        chunk_size: int = 6000,
        generate_questions: bool = True
    ) -> dict:
        """
        FAST ingestion pipeline using PyMuPDF (no Unstructured).

        Phases:
        1. Extract text with PyMuPDF (fast, ~2s)
        2. Extract images with PyMuPDF (fast, ~1s)
        3. Upload images to R2
        4. Create hierarchy with LLM (Groq Llama 3.1 8B, ~1s)
        5. Match images to chunks
        6. Generate embeddings + questions (Groq GPT-OSS-120B, ~5s parallel)
        7. Persist to Neo4j

        This is ~10x faster than the Unstructured-based pipeline.

        Args:
            file_path: Path to PDF file
            doc_id: Unique document identifier
            user_id: User who uploaded the document
            title: Optional document title
            chunk_size: Target characters per chunk
            generate_questions: Whether to generate questions

        Returns:
            Summary dict with counts and timing
        """
        from pathlib import Path

        total_start = time.time()
        print(f"\n{'='*60}")
        print(f"🚀 FAST Ingestion Pipeline (PyMuPDF)")
        print(f"📄 Document: {Path(file_path).name}")
        print(f"🆔 Doc ID: {doc_id}")
        print(f"{'='*60}\n")

        # Default title from filename
        if not title:
            title = Path(file_path).stem

        # ===== Phase 1: Fast text extraction with OCR fallback for problem pages =====
        # Uses PyMuPDF for fast extraction, with encoding detection and
        # DeepInfra Mistral-Small OCR for pages with LaTeX/equation issues
        from lib.encoding_check import extract_pdf_with_fallback
        from lib.text_chunker import chunk_pages

        print("⏳ Phase 1: Text Extraction (PyMuPDF + OCR fallback)...")
        pages_text = extract_pdf_with_fallback(file_path, doc_id=doc_id)

        # Chunk the pages with smart boundaries
        chunks = chunk_pages(pages_text, max_chars=chunk_size, min_chars=500, overlap=200)

        # Add chunk_index to match expected format
        for i, chunk in enumerate(chunks):
            chunk["chunk_index"] = i

        print(f"✅ Phase 1 complete: {len(chunks)} chunks from {len(pages_text)} pages\n")

        # ===== Phase 2: Image extraction + Vision filtering =====
        print("⏳ Phase 2: Image Extraction (PyMuPDF)...")
        images_by_page = extract_images_from_pdf(file_path)
        image_index = {}
        images_before_vision = sum(len(imgs) for imgs in images_by_page.values()) if images_by_page else 0

        if images_by_page:
            # Apply vision LLM filtering to keep only educational images
            print("⏳ Phase 2b: Vision LLM Classification...")
            images_by_page = filter_and_describe_images(images_by_page)
            images_after_vision = sum(len(imgs) for imgs in images_by_page.values())

            if images_by_page:
                image_index = upload_images_to_r2(images_by_page, doc_id)
                print(f"✅ Phase 2 complete: {images_after_vision}/{images_before_vision} educational images uploaded\n")
            else:
                print(f"✅ Phase 2 complete: 0/{images_before_vision} images passed vision filter\n")
        else:
            print("✅ Phase 2 complete: No images found\n")

        # ===== Phase 3: Create hierarchy with LLM =====
        print("⏳ Phase 3: Hierarchy Creation (LLM)...")
        # Create temporary ContentBlock-like objects for hierarchy function
        temp_blocks = []
        for chunk in chunks:
            block = ContentBlock()
            block.text_content = chunk["text"]
            block.page_number = chunk["page_start"]
            block.page_start = chunk["page_start"]
            block.page_end = chunk["page_end"]
            block.chunk_index = chunk["chunk_index"]
            temp_blocks.append(block)

        hierarchy = create_hierarchy_with_llm(temp_blocks)
        temp_blocks = apply_hierarchy_to_chunks(temp_blocks, hierarchy)
        print(f"✅ Phase 3 complete: Hierarchy applied\n")

        # ===== Phase 4: Match images to chunks =====
        if images_by_page and image_index:
            print("⏳ Phase 4: Image-Chunk Matching...")
            # Images already have descriptions from vision filtering
            temp_blocks = match_images_to_chunks(
                temp_blocks,
                images_by_page,
                image_index,
                vision_llm=None  # Descriptions already in images_by_page from filter_and_describe_images
            )
            print(f"✅ Phase 4 complete\n")
        else:
            print("⏳ Phase 4: Skipped (no images)\n")

        # ===== Phase 5: Embeddings + Questions (parallel) =====
        if generate_questions:
            print("⏳ Phase 5: Embeddings + Questions (parallel)...")
            # Build combined context WITH figure descriptions for QP generation
            for block in temp_blocks:
                block.combined_context = build_combined_context_with_figures(block)

            temp_blocks = self.enrich_content_blocks(temp_blocks, parallel_questions=True)
            # Note: image linking now happens inside _generate_questions via figure_ref
            print(f"✅ Phase 5 complete\n")
        else:
            # Just generate embeddings
            print("⏳ Phase 5: Embeddings only...")
            for block in temp_blocks:
                block.combined_context = build_combined_context_with_figures(block)
                if block.combined_context:
                    block.embeddings = embed_text(block.combined_context)
            print(f"✅ Phase 5 complete\n")

        # ===== Phase 6: Persist to Neo4j =====
        print("⏳ Phase 6: Neo4j Persistence...")
        doc_meta = {
            "user_id": user_id,
            "doc_id": doc_id,
            "title": title,
            "source": Path(file_path).name
        }
        self.persist_to_neo4j(doc_id, doc_meta, temp_blocks)
        print(f"✅ Phase 6 complete\n")

        # ===== Summary =====
        total_elapsed = time.time() - total_start
        total_questions = sum(len(b.questions) for b in temp_blocks)
        total_images = sum(len(b.image_urls) for b in temp_blocks)

        chapters = set(b.chapter_title for b in temp_blocks if b.chapter_title)
        sections = set(b.section_title for b in temp_blocks if b.section_title)

        summary = {
            "doc_id": doc_id,
            "title": title,
            "content_blocks": len(temp_blocks),
            "chapters": len(chapters),
            "sections": len(sections),
            "questions": total_questions,
            "images_extracted": sum(len(imgs) for imgs in images_by_page.values()) if images_by_page else 0,
            "images_uploaded": len(image_index),
            "images_matched": total_images,
            "elapsed_seconds": round(total_elapsed, 2),
            "pipeline": "fast"
        }

        print(f"{'='*60}")
        print(f"✅ FAST Ingestion Complete!")
        print(f"   📦 Blocks: {summary['content_blocks']}")
        print(f"   📖 Chapters: {summary['chapters']}, Sections: {summary['sections']}")
        print(f"   ❓ Questions: {summary['questions']}")
        print(f"   🖼️  Images: {summary['images_extracted']} extracted, {summary['images_matched']} matched")
        print(f"   ⏱️  Time: {total_elapsed:.2f}s ({total_elapsed/60:.1f} min)")
        print(f"{'='*60}\n")

        return summary


if __name__ == "__main__":
    import sys
    import uuid

    # Default test file or command line arg
    test_file = sys.argv[1] if len(sys.argv) > 1 else "chapter1.pdf"
    use_fast = "--fast" in sys.argv or "-f" in sys.argv

    config = {
        "vision_llm": "gpt-4o-mini",
        "text_llm": "gpt-4.1"
    }

    pipeline = IngestionPipeline(config)

    # Generate unique doc_id
    doc_id = f"doc_{uuid.uuid4().hex[:8]}"

    if use_fast:
        # Run the FAST pipeline (PyMuPDF, no Unstructured)
        print("🏃 Using FAST pipeline (PyMuPDF)")
        result = pipeline.ingest_document_fast(
            file_path=test_file,
            doc_id=doc_id,
            user_id="test_user",
            generate_questions=True
        )
    else:
        # Run the standard pipeline (Unstructured)
        print("📚 Using STANDARD pipeline (Unstructured)")
        result = pipeline.ingest_document(
            file_path=test_file,
            doc_id=doc_id,
            user_id="test_user",
            extract_images=True,
            create_hierarchy=True,
            generate_questions=True
        )

    print(f"\n📊 Result: {json.dumps(result, indent=2)}")
